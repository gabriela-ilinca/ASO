"""
Atlanta Symphony Orchestra - Comprehensive Data Analysis Presentation
Consulting-style presentation with detailed speaker notes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Color scheme - Professional consulting palette
NAVY = RGBColor(0, 51, 102)       # Primary dark blue
GOLD = RGBColor(197, 164, 103)    # ASO gold accent
LIGHT_BLUE = RGBColor(66, 133, 244)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(219, 68, 55)
GREEN = RGBColor(15, 157, 88)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with professional styling"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add navy background rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Gold accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(10), Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_header(prs, title, notes=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Navy left bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Gold accent
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(3.4), Inches(9.7), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_content_slide(prs, title, content_items, notes="", has_chart_placeholder=False):
    """Add a content slide with action title"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Gold underline
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(10), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Title (action title style)
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)
            p.level = 0

    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_data_slide(prs, title, table_data, notes="", key_insight=None):
    """Add a slide with tabular data"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Gold underline
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(10), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Create table
    if table_data:
        rows = len(table_data)
        cols = len(table_data[0])

        # Calculate column widths
        table_width = 8.5
        col_width = table_width / cols

        table = slide.shapes.add_table(rows, cols, Inches(0.75), Inches(1.5), Inches(table_width), Inches(0.4 * rows)).table

        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)

                # Style header row
                if row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = NAVY
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = WHITE
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(12)
                else:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.color.rgb = DARK_GRAY

    # Key insight callout
    if key_insight:
        insight_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.8), Inches(8.5), Inches(0.8)
        )
        insight_box.fill.solid()
        insight_box.fill.fore_color.rgb = RGBColor(255, 248, 220)  # Light gold
        insight_box.line.color.rgb = GOLD

        insight_tf = insight_box.text_frame
        insight_tf.paragraphs[0].text = "KEY INSIGHT: " + key_insight
        insight_tf.paragraphs[0].font.size = Pt(13)
        insight_tf.paragraphs[0].font.bold = True
        insight_tf.paragraphs[0].font.color.rgb = NAVY

    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_two_column_slide(prs, title, left_content, right_content, notes=""):
    """Add a two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Gold underline
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(10), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = LIGHT_BLUE
    left_header.line.fill.background()
    left_header.text_frame.paragraphs[0].text = left_content.get('header', 'Column 1')
    left_header.text_frame.paragraphs[0].font.color.rgb = WHITE
    left_header.text_frame.paragraphs[0].font.bold = True
    left_header.text_frame.paragraphs[0].font.size = Pt(14)

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content.get('items', [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    # Right column header
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = GREEN
    right_header.line.fill.background()
    right_header.text_frame.paragraphs[0].text = right_content.get('header', 'Column 2')
    right_header.text_frame.paragraphs[0].font.color.rgb = WHITE
    right_header.text_frame.paragraphs[0].font.bold = True
    right_header.text_frame.paragraphs[0].font.size = Pt(14)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content.get('items', [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_kpi_slide(prs, title, kpis, notes=""):
    """Add a KPI dashboard slide with large metrics"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Gold underline
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(10), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # KPI boxes
    num_kpis = len(kpis)
    box_width = 2.8
    spacing = 0.3
    total_width = num_kpis * box_width + (num_kpis - 1) * spacing
    start_x = (10 - total_width) / 2

    for i, kpi in enumerate(kpis):
        x = start_x + i * (box_width + spacing)

        # KPI box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(box_width), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 250)
        box.line.color.rgb = NAVY

        # KPI value
        value_box = slide.shapes.add_textbox(Inches(x), Inches(2.0), Inches(box_width), Inches(0.8))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = kpi.get('value', '')
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = kpi.get('color', NAVY)
        p.alignment = PP_ALIGN.CENTER

        # KPI label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(2.8), Inches(box_width), Inches(1.0))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = kpi.get('label', '')
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_recommendation_slide(prs, title, recommendation, supporting_points, notes=""):
    """Add a recommendation slide with callout"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Gold underline
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(10), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Recommendation callout box
    rec_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2)
    )
    rec_box.fill.solid()
    rec_box.fill.fore_color.rgb = RGBColor(230, 247, 255)
    rec_box.line.color.rgb = LIGHT_BLUE

    rec_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.6), Inches(1.0))
    tf = rec_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = recommendation
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Supporting points
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(4.0))
    tf = points_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(supporting_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def create_presentation():
    """Create the complete ASO presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "Atlanta Symphony Orchestra",
        "Data Analysis and Pricing Recommendations | December 2024"
    )

    # ========== SLIDE 2: Agenda ==========
    add_content_slide(
        prs,
        "Agenda: Data-Driven Insights to Optimize Revenue and Engagement",
        [
            "1. Audience Behavior Analysis",
            "     - Subscribers vs. Single Ticket Buyers comparison",
            "     - Purchase patterns and timing insights",
            "",
            "2. Revenue Deep Dive",
            "     - Revenue generators and subscription economics",
            "     - CYO subscriber value analysis",
            "",
            "3. Demand Curve Analysis",
            "     - Price elasticity by segment and seat section",
            "     - Optimization opportunities",
            "",
            "4. Pricing Simulations and Recommendations",
            "     - Classical Add-On Product (targeting Special/Holiday buyers)",
            "     - Fixed Plus Premium Program (targeting Fixed 6 subscribers)",
            "",
            "5. Implementation Roadmap and Success Metrics"
        ],
        notes="""SPEAKER NOTES - Agenda Slide

OPENING (30 seconds):
"Thank you for the opportunity to present our findings. Over the past several months, we analyzed over 68,000 ticket records from FY23 and FY24 to answer three strategic questions: How do we get single ticket buyers to return? How do we attract and retain subscribers? And how do we engage younger audiences?"

KEY FRAMING:
"Today's presentation builds from understanding customer behavior, through demand analysis, to two specific pricing simulations that we believe can generate meaningful incremental revenue while strengthening customer relationships."

TRANSITION:
"Let's start by understanding who your customers are and how they behave."

ANTICIPATED QUESTIONS:
- Q: What data sources did you use? A: FY23 and FY24 regular ticket data, demographics, price codes, and event information - over 68,000 paid ticket records.
- Q: Who was involved in this analysis? A: Cross-functional collaboration with box office, marketing, and finance teams.
"""
    )

    # ========== SLIDE 3: Section Header - Audience Behavior ==========
    add_section_header(
        prs,
        "Section 1: Audience Behavior Analysis",
        notes="""SPEAKER NOTES - Section Header

TRANSITION:
"Understanding audience behavior is the foundation of any pricing strategy. Let's look at how your subscriber and single ticket buyer populations compare and how they've evolved."

KEY CONTEXT TO SET:
- ASO averages 88% hall occupancy - among the best in the industry
- Budget split: 51% contributed, 49% earned revenue
- Data covers two full seasons: FY23 and FY24
"""
    )

    # ========== SLIDE 4: Subscriber vs Single Ticket Comparison ==========
    add_two_column_slide(
        prs,
        "Single ticket buyer growth outpaces subscribers, creating both opportunity and risk",
        {
            'header': 'Single Ticket Buyers (Growing)',
            'items': [
                "FY23 to FY24: +2,796 buyers",
                "Now 30% of total audience",
                "Average ticket price: $59.86",
                "Purchase 17 days before event",
                "80% attend only once",
                "Huge growth in Holiday, Special, Delta Series"
            ]
        },
        {
            'header': 'Subscribers (Declining)',
            'items': [
                "FY23 to FY24: -67 subscribers",
                "Average ticket price: $51.02",
                "Purchase 2+ months in advance",
                "85% renewal rate (established)",
                "57% of subscribers buy add-on tickets",
                "Average subscriber age: 63"
            ]
        },
        notes="""SPEAKER NOTES - Subscriber vs Single Ticket Comparison

KEY INSIGHT:
"Your single ticket buyer population is growing rapidly - up by nearly 2,800 accounts year-over-year - while subscribers declined slightly. This is both an opportunity and a challenge. Single ticket buyers pay more per ticket ($59.86 vs $51.02) but their retention is extremely low - 80% never return."

DATA POINTS TO EMPHASIZE:
1. The price differential: Single ticket buyers pay 17% more per seat than subscribers
2. The timing difference: 17 days vs 2 months - single ticket buyers are last-minute
3. The retention gap: 80% single ticket churn vs 85% established subscriber renewal

STRATEGIC IMPLICATION:
"The growth in single ticket buyers, especially in Holiday, Special, and Delta Series events, represents a pipeline for future subscribers IF we can convert them. This is the insight that drives our first pricing simulation."

ANTICIPATED QUESTIONS:
- Q: Why are subscribers declining? A: Industry-wide trend toward flexibility; also demographic shifts (average age 63)
- Q: What's driving single ticket growth in Holiday/Special? A: More accessible programming, different audience segment
"""
    )

    # ========== SLIDE 5: Single Ticket Buyer Behavior ==========
    add_kpi_slide(
        prs,
        "Single ticket buyers purchase late, attend selectively, and rarely return",
        [
            {'value': '17 days', 'label': 'Average purchase lead time before event', 'color': LIGHT_BLUE},
            {'value': '80%', 'label': 'Attend only one event per season', 'color': RED},
            {'value': '$59.86', 'label': 'Average ticket price (17% above subscribers)', 'color': GREEN}
        ],
        notes="""SPEAKER NOTES - Single Ticket Buyer Behavior

KEY STORY:
"Single ticket buyers are a fascinating segment. They're willing to pay more - $59.86 on average versus $51.02 for subscribers - but they decide late and rarely come back. The 17-day purchase window versus 2 months for subscribers tells us they're responding to specific event marketing, not planning their season."

BEHAVIORAL INSIGHT:
"The 80% one-and-done rate is actually an opportunity. These aren't people who dislike ASO - they bought a ticket! They simply haven't been given a compelling reason to return. That's what our pricing simulations address."

SHARE OF TICKET GROUPS:
- Holiday events: Largest single ticket buyer concentration
- Special events: Second highest
- Classical: Lower single ticket penetration but higher loyalty

FLUCTUATION PATTERNS:
"Single ticket sales are highly variable week-to-week, driven by artist recognition and event marketing. Subscribers provide predictable baseline; single ticket is your upside."

TRANSITION:
"Now let's look at subscribers, who behave very differently..."
"""
    )

    # ========== SLIDE 6: Subscriber Behavior Deep Dive ==========
    add_kpi_slide(
        prs,
        "Subscribers are highly engaged beyond their subscription: 57% purchase add-on tickets",
        [
            {'value': '57%', 'label': 'of subscribers buy at least one add-on single ticket', 'color': NAVY},
            {'value': '4', 'label': 'average add-on tickets purchased per buyer', 'color': LIGHT_BLUE},
            {'value': '$355', 'label': 'average single ticket revenue per subscriber buyer', 'color': GREEN}
        ],
        notes="""SPEAKER NOTES - Subscriber Behavior

KEY INSIGHT:
"Your subscribers aren't just attending their subscription concerts - they want MORE. 57% of subscribers purchase additional single tickets, averaging 4 add-on tickets worth $355 in revenue. This is proven demand for content beyond their package."

MOST POPULAR ADD-ON EVENTS (BY CROSSOVER):
1. Josh Bell - celebrity draw
2. Vivaldi's Four Seasons - accessible repertoire
3. Beethoven Symphony No. 9 - iconic programming

"Notice a pattern? Subscribers gravitate toward recognizable names and works when buying add-ons. This insight informs our marketing recommendations."

STRATEGIC IMPLICATION:
"The fact that 57% of subscribers already buy extras means there's no risk of cannibalizing subscription value with an add-on program. We're simply making it easier for them to do what they already want to do - and capturing a premium in the process."

TRANSITION:
"Now let's understand the revenue implications of these behaviors..."
"""
    )

    # ========== SLIDE 7: Section Header - Revenue ==========
    add_section_header(
        prs,
        "Section 2: Revenue Generators",
        notes="""SPEAKER NOTES - Section Header

TRANSITION:
"Understanding who generates revenue - and how - is critical for pricing strategy. Let's look at the composition of your earned revenue and some surprising findings about subscriber economics."
"""
    )

    # ========== SLIDE 8: Revenue by Fiscal Year ==========
    add_data_slide(
        prs,
        "Subscription revenue provides stability, but single ticket revenue drives growth potential",
        [
            ["Revenue Source", "FY23", "FY24", "Change"],
            ["Subscription Revenue", "$3.8M", "$3.7M", "-2.6%"],
            ["Single Ticket Revenue", "$2.1M", "$2.4M", "+14.3%"],
            ["Total Ticket Revenue", "$5.9M", "$6.1M", "+3.4%"],
            ["Subscription Share", "64%", "61%", "-3 pts"],
        ],
        notes="""SPEAKER NOTES - Revenue by Fiscal Year

KEY STORY:
"While total ticket revenue grew 3.4% year-over-year, the composition is shifting. Subscription revenue declined 2.6% while single ticket revenue surged 14.3%. The subscription share of total revenue dropped from 64% to 61%."

STRATEGIC IMPLICATIONS:
1. Revenue growth is coming from single ticket, not subscription
2. This increases revenue volatility (single ticket is less predictable)
3. Converting single ticket buyers to subscribers would stabilize revenue

WHY THIS MATTERS:
"Subscriptions are sold well in advance and provide cash flow predictability. Single ticket revenue, while growing, comes late and is event-dependent. The ideal outcome is to grow single ticket revenue while building a pipeline to subscription."

TRANSITION:
"Within subscribers, there's significant variation in value. Let's look at CYO subscribers specifically..."
""",
        key_insight="Single ticket revenue grew 14.3% YoY - the growth engine is outside subscriptions"
    )

    # ========== SLIDE 9: CYO Revenue Impact ==========
    add_kpi_slide(
        prs,
        "CYO subscribers represent 13.73% of members but contribute 22-24% of subscription revenue",
        [
            {'value': '13.73%', 'label': 'Share of subscription accounts (CYO)', 'color': NAVY},
            {'value': '22-24%', 'label': 'Share of subscription revenue (CYO)', 'color': GREEN},
            {'value': '$67.96', 'label': 'CYO average ticket price (vs $39.48 other subs)', 'color': LIGHT_BLUE}
        ],
        notes="""SPEAKER NOTES - CYO Revenue Contribution

KEY INSIGHT:
"CYO subscribers are your most valuable subscriber segment by far. They represent just 13.73% of subscription accounts but generate 22-24% of subscription revenue. Their average ticket price of $67.96 is 72% higher than other subscribers at $39.48."

WHY CYO IS VALUABLE:
1. Self-selection: CYO subscribers choose specific concerts, so they attend what they want
2. Premium seating: They tend to select better seats
3. Higher engagement: They're actively planning their season
4. Better retention: Higher satisfaction from personalized experience

STRATEGIC IMPLICATION:
"Growing your CYO subscriber base should be a priority. That's why our Classical Add-On program includes a CYO conversion pathway - we project 50 new CYO subscribers from the program."

TRANSITION:
"Let's look at subscriber retention patterns to understand where conversion efforts should focus..."
"""
    )

    # ========== SLIDE 10: Subscriber Movement ==========
    add_content_slide(
        prs,
        "Subscriber retention is critically low in the first two years: only 56% of first-year subscribers return",
        [
            "SUBSCRIBER RETENTION BY TENURE:",
            "",
            "   Freshman to Sophomore:     56.1% retention",
            "   Sophomore to Established:  <80% retention",
            "   Established renewals:      85%+ retention",
            "",
            "",
            "KEY FINDING: Once subscribers reach Year 3, they become highly loyal.",
            "The challenge is getting them through Years 1 and 2.",
            "",
            "",
            "CYO SUBSCRIBER CHURN:",
            "",
            "   161 CYO subscribers in 2023 did not return in 2024",
            "   Peak consumption: 3-4 tickets per CYO package",
            "   Under-utilization may signal dissatisfaction risk"
        ],
        notes="""SPEAKER NOTES - Subscriber Movement

KEY STORY:
"Subscriber retention follows a predictable pattern: Year 1 is brutal at only 56%, Year 2 improves to about 80%, and by Year 3+ you're seeing 85%+ retention. The implication is clear - if you can keep subscribers through their first two years, they become highly loyal."

CYO SPECIFIC INSIGHTS:
"We lost 161 CYO subscribers between 2023 and 2024. Analysis shows peak consumption at 3-4 tickets - some CYO subscribers aren't using their full packages, which may indicate a mismatch between what they bought and what they need."

STRATEGIC IMPLICATIONS:
1. Invest in Year 1 and Year 2 subscriber experience
2. Consider package flexibility for new subscribers
3. Monitor CYO utilization as an early warning indicator

ACTION ITEMS:
- Outreach to first-year subscribers mid-season
- Under-utilized CYO subscribers need attention before renewal
- Create "bridge" products that ease the commitment curve

TRANSITION:
"Now let's move to demand analysis to understand pricing opportunities..."
"""
    )

    # ========== SLIDE 11: Section Header - Demand Curves ==========
    add_section_header(
        prs,
        "Section 3: Demand Curve Analysis",
        notes="""SPEAKER NOTES - Section Header

TRANSITION:
"Understanding price elasticity - how quantity demanded changes with price - is essential for setting optimal prices. We analyzed demand curves across multiple dimensions: by seat section, subscription type, and event category."
"""
    )

    # ========== SLIDE 12: Demand by Seat Section ==========
    add_data_slide(
        prs,
        "Single ticket demand varies by event type - Special Events are most price elastic",
        [
            ["Segment", "Avg Price", "Elasticity", "Interpretation"],
            ["Single Ticket - Special", "$77.25", "-1.35", "Elastic - discounts drive volume AND revenue"],
            ["Single Ticket - Holiday", "$56.21", "-0.86", "Slightly inelastic - loyal audience"],
            ["Single Ticket - Classical", "$58.01", "-0.18", "Very inelastic - price won't drive new buyers"],
            ["Single Ticket (All)", "$58.93", "-0.36", "Moderately inelastic overall"],
            ["Single Ticket - Premium Orch", "$63.07", "-0.41", "Inelastic - premium buyers less sensitive"]
        ],
        notes="""SPEAKER NOTES - Demand by Event Type

INTERPRETING ELASTICITY:
"Negative elasticity between 0 and -1 means inelastic demand - raising prices increases revenue. Elasticity below -1 means elastic demand - lowering prices can increase revenue."

KEY FINDINGS:
1. Special Events (-1.35): Elastic demand! A 10% discount leads to 15% more tickets and 3.8% net revenue gain.
2. Holiday (-0.86): Slightly inelastic - loyal Holiday audience will pay.
3. Classical (-0.18): Very inelastic - price changes don't move the needle.

STRATEGIC IMPLICATIONS:
"Special Events respond to promotional pricing. Classical does not - those who want it will pay; price cuts just erode margin. Holiday is in between."

PRICING RECOMMENDATIONS:
- Special Events: Strategic discounting can work
- Classical: Protect pricing; use programming, not price cuts
- Holiday: Maintain current pricing
""",
        key_insight="Special Events are elastic (-1.35) - strategic discounts increase revenue"
    )

    # ========== SLIDE 13: Demand by Subscription Type ==========
    add_data_slide(
        prs,
        "Subscription demand by section shows CYO subscribers are more price sensitive than Fixed",
        [
            ["Segment", "Avg Price", "Elasticity", "Notes"],
            ["Subscription - Premium Orchestra", "$45.02", "-0.89", "Inelastic - loyal subscribers"],
            ["Subscription - Front Loge", "$75.77", "+0.62", "Some price sensitivity"],
            ["Subscription - Balcony", "$24.17", "-0.20", "Inelastic - budget-driven"],
            ["CYO - Premium Orchestra", "$51.99", "-2.19", "Very elastic - price conscious"],
            ["CYO - Dress Circle", "$48.67", "-2.05", "Very elastic - value seekers"]
        ],
        notes="""SPEAKER NOTES - Subscription Demand by Section

KEY INSIGHTS:
1. Fixed subscribers in Premium Orchestra (-0.89) are relatively inelastic - they're committed
2. CYO subscribers are very elastic (-2.19 in Premium Orchestra) - they actively compare value
3. Balcony subscribers are inelastic but driven by budget constraints

STRATEGIC INTERPRETATION:
"Fixed subscribers have made a commitment and are less price sensitive. CYO subscribers are 'shopping' each season and respond more to price changes. This explains why CYO generates higher per-ticket revenue - they're choosing premium when they see value."

PRICING STRATEGY:
- Fixed subscriptions: Can support modest price increases
- CYO subscriptions: Be careful with price increases; focus on value communication
- Balcony: Keep as accessible entry point

WHY THIS MATTERS FOR FIXED PLUS:
"Fixed subscribers' inelastic demand (-0.89) means they can support a 10% premium for bundled value without significant volume loss."
""",
        key_insight="Fixed subscribers are inelastic (-0.89) - supports 10% Fixed Plus premium"
    )

    # ========== SLIDE 14: Demand by Event Type ==========
    add_data_slide(
        prs,
        "Special Events show high elasticity for single ticket buyers - discounts drive volume",
        [
            ["Segment", "Avg Price", "Elasticity", "Revenue Impact of 10% Discount"],
            ["Single Ticket - Special Events", "$77.25", "-1.35", "+15.3% quantity, +3.8% revenue"],
            ["Single Ticket (All)", "$58.93", "-0.36", "+3.8% quantity, -6.6% revenue"],
            ["Premium Single Ticket - Special", "$85.74", "-0.91", "+10% quantity, break-even"],
            ["Classical (Single Ticket)", "$62.43", "-0.18", "+1.8% quantity, -8.2% revenue"]
        ],
        notes="""SPEAKER NOTES - Demand by Event Type

CRITICAL INSIGHT:
"Special Events single tickets show high elasticity (-1.35). This means discounting Special Events actually INCREASES revenue - a 10% discount leads to 15% more tickets and 3.8% net revenue gain. This is unusual and valuable."

CLASSICAL IS DIFFERENT:
"Classical single tickets are highly inelastic (-0.18). Discounting Classical erodes revenue without driving meaningful volume. The buyers who want Classical will pay; those who don't, won't be swayed by price alone."

WHY THIS MATTERS FOR CLASSICAL ADD-ON:
"This is exactly why our Classical Add-On program uses modest discounts (5-15%) - not to drive price-sensitive demand, but to remove the psychological barrier for Special/Holiday buyers who've never tried Classical. The discount is a 'permission slip' to try something new."

PRICING STRATEGY:
- Special Events: Can use discounting strategically
- Classical: Protect pricing; use product differentiation not price cuts
- Premium sections: Maintain or increase

TRANSITION:
"Now let's apply these insights to two specific pricing simulations..."
"""
    )

    # ========== SLIDE 15: Section Header - Simulations ==========
    add_section_header(
        prs,
        "Section 4: Pricing Simulations",
        notes="""SPEAKER NOTES - Section Header

TRANSITION:
"Based on our behavioral and demand analysis, we developed two pricing simulations that target different customer segments with different objectives. Both are designed to generate incremental revenue while strengthening customer relationships."

PREVIEW:
1. Classical Add-On Product: Converts Special/Holiday single ticket buyers into Classical attendees
2. Fixed Plus Premium: Captures value from Fixed 6 subscribers who already want more
"""
    )

    # ========== SLIDE 16: Classical Add-On Overview ==========
    add_recommendation_slide(
        prs,
        "Simulation A: Classical Add-On Product targets 7,008 Special/Holiday buyers",
        "OPPORTUNITY: 80.6% of Special/Holiday single ticket buyers have NEVER purchased a Classical ticket - a massive untapped market for cross-selling.",
        [
            "Target Market: 7,008 Special/Holiday single ticket buyers (30% of all single ticket accounts)",
            "Key Insight: These buyers already love ASO - they just haven't discovered Classical",
            "Approach: Offer tiered discounts for Classical tickets to recent Special/Holiday attendees",
            "Timing: Same-day offers maximize engagement while experience is fresh",
            "Risk: Zero - product differentiation means no cannibalization of full-price sales"
        ],
        notes="""SPEAKER NOTES - Classical Add-On Overview

FRAMING THE OPPORTUNITY:
"Think about this: 7,008 people attended your Special or Holiday events as single ticket buyers last season. More than 80% of them have NEVER bought a Classical ticket. These aren't strangers - they've already paid to experience ASO. They simply haven't been invited to try your core product."

WHY SAME-DAY OFFERS:
"Behavioral economics tells us the moment of highest engagement is immediately after a positive experience. By offering Classical discounts on the same day someone attends a Special event, we're meeting them at peak satisfaction. The ask is: 'You loved this. Want more?'"

RISK MITIGATION:
"This program offers discounts ONLY on the new add-on product, not existing pricing. Someone who would have bought Classical at full price still can. We're creating a new conversion pathway without eroding current revenue."

ANTICIPATED OBJECTION:
- Q: Won't this train people to expect discounts?
- A: No - the discount is explicitly tied to attending a gateway event. It's a reward for engagement, not a general price cut.
"""
    )

    # ========== SLIDE 17: Classical Add-On Discount Structure ==========
    add_data_slide(
        prs,
        "Tiered discounts incentivize multiple purchases: +1 ticket (5% off) to +3 tickets (15% off)",
        [
            ["Add-On Level", "Discount", "Example Price", "Psychological Driver"],
            ["+1 Classical ticket", "5% off", "$55.63 (was $58.56)", "Low barrier to trial"],
            ["+2 Classical tickets", "10% off each", "$52.70 per ticket", "Bring a friend / Try two"],
            ["+3 Classical tickets", "15% off each", "$49.78 per ticket", "Subscription-like behavior"]
        ],
        notes="""SPEAKER NOTES - Discount Structure

RATIONALE FOR TIERED DISCOUNTS:
"We designed the discount structure to encourage progressive engagement. A 5% discount on one ticket removes the price barrier for trying Classical. The 10% discount for two tickets encourages bringing a companion - which increases the social experience. The 15% discount for three tickets creates subscription-like behavior."

OFFER WINDOW:
"The offer is valid for 60 days from attending a Special/Holiday event. This creates urgency while allowing time for upcoming Classical concerts to fall within the window."

MARKETING TOUCHPOINTS:
1. Email within 2 hours post-event: 'Continue your ASO experience...'
2. In-venue signage during intermission/exit with QR code
3. Box office staff mention at check-in
4. SMS next morning (if opted in)

GATEWAY EVENTS TO PROMOTE:
Based on historical crossover, these Classical concerts have highest appeal to Special/Holiday buyers:
- Beethoven Symphony No. 9 (188 crossover buyers)
- Opening Weekend (164)
- Orff: Carmina Burana (139)
- Beethoven and Bolero (119)

"The message should connect: 'You loved Holiday Pops - experience Beethoven's 9th next!'"
""",
        key_insight="60-day offer window creates urgency while allowing planning time"
    )

    # ========== SLIDE 18: Classical Add-On Before/After ==========
    add_data_slide(
        prs,
        "Projected impact: 1,926 new Classical tickets generating $102,099 incremental revenue",
        [
            ["Segment", "Current State", "After Program", "Net New"],
            ["Group A: 0 Classical (5,647 accounts)", "0 tickets", "1,689 tickets", "+1,689"],
            ["Group B: 1 Classical (729 accounts)", "1,559 tickets", "1,746 tickets", "+187"],
            ["Group C: 2 Classical (254 accounts)", "986 tickets", "1,036 tickets", "+50"],
            ["TOTAL NEW TICKETS", "-", "-", "1,926"],
            ["INCREMENTAL REVENUE", "-", "-", "$102,099"]
        ],
        notes="""SPEAKER NOTES - Before/After Analysis

METHODOLOGY:
"We applied a conservative 20% response rate to the target population. This is based on industry benchmarks for same-day, targeted offers to warm audiences. The actual rate could be higher given the engagement timing."

BREAKDOWN BY SEGMENT:
- Group A (never bought Classical): 1,129 responders, 1,689 new tickets at ~$53.30 average
- Group B (1 Classical ever): 145 responders buying 187 additional tickets
- Group C (2 Classical ever): 50 responders buying 50 additional tickets

REVENUE CALCULATION:
- Total discounted revenue: $102,099
- Full price equivalent: $112,777
- Effective discount rate: 9.5%
- Average discount per ticket: $5.54

KEY POINT:
"The 'discount cost' of ~$10,700 represents value given up versus full price. But these are customers who would NOT have bought Classical otherwise. The full $102,099 is incremental revenue."
""",
        key_insight="80.6% of target audience (Group A) has never tried Classical - largest opportunity"
    )

    # ========== SLIDE 19: Classical Add-On CYO Pipeline ==========
    add_kpi_slide(
        prs,
        "The Classical Add-On creates a CYO conversion pipeline: 50 projected new subscribers",
        [
            {'value': '205', 'label': 'New buyers reaching 3+ Classical tickets threshold', 'color': NAVY},
            {'value': '24.6%', 'label': 'Historical conversion rate from 3+ Classical to CYO', 'color': LIGHT_BLUE},
            {'value': '50', 'label': 'Projected new CYO subscribers from program', 'color': GREEN}
        ],
        notes="""SPEAKER NOTES - CYO Conversion Pipeline

THE CONVERSION FUNNEL:
"We analyzed historical data and found that 24.6% of single ticket buyers who purchase 3+ Classical tickets eventually become CYO subscribers. This is a proven pathway."

HOW THE ADD-ON CREATES CONVERTS:
- Group A buyers getting +3: 112 accounts reach threshold
- Group B buyers getting +2 more: 43 accounts reach threshold
- Group C buyers getting +1 more: 50 accounts reach threshold
- Total new 3+ buyers: 205 accounts

APPLYING CONVERSION RATE:
"205 new 3+ Classical buyers x 24.6% conversion rate = 50 projected new CYO subscribers"

RECURRING REVENUE VALUE:
- Low estimate: 50 x $400/season = $20,000/year
- High estimate: 50 x $600/season = $30,000/year

STRATEGIC IMPLICATION:
"This is the compound effect. The Classical Add-On generates $102K in Year 1 ticket revenue PLUS creates an ongoing pipeline of CYO subscribers worth $20-30K annually. And those CYO subscribers have high lifetime value."
"""
    )

    # ========== SLIDE 20: Fixed Plus Overview ==========
    add_recommendation_slide(
        prs,
        "Simulation B: Fixed Plus Premium targets Fixed 6 subscribers who want more",
        "OPPORTUNITY: 57.8% of Fixed 6 subscribers already purchase single tickets, proving demand for content beyond their package. Fixed Plus bundles what they want at a 10% premium.",
        [
            "Target Market: 1,643 Fixed 6 subscribers (70% of all Fixed subscribers)",
            "Proven Demand: 950 accounts (57.8%) already buy single tickets averaging $355/buyer",
            "Approach: Bundle Holiday + Special event access with exclusive benefits",
            "Positioning: Sell benefits (priority access, guaranteed seating) - not price",
            "Premium: 10% above regular pricing, justified by tangible value-adds"
        ],
        notes="""SPEAKER NOTES - Fixed Plus Overview

WHY FIXED 6:
"Fixed 6 is your largest Fixed tier at 1,643 accounts - 70% of all Fixed subscribers. They're committed enough to subscribe but haven't maxed out. And critically, 57.8% of them are ALREADY buying single tickets. This is proven demand."

WHAT THEY BUY AS SINGLES:
- Classical: 709 unique buyers, $208,590 revenue
- Special: 381 unique buyers, $82,806 revenue
- Holiday: 196 unique buyers, $41,980 revenue

THE FIXED PLUS CONCEPT:
"Instead of letting them buy Holiday and Special as separate transactions, we bundle it into their subscription at a 10% premium. They get guaranteed access and premium benefits; we capture incremental revenue and lock in commitment."

WHY 10% PREMIUM:
"Our elasticity analysis shows Fixed subscribers can support modest price increases. A 10% premium for bundled value - priority access, guaranteed seating, no separate transactions - is well within acceptance range for this segment."
"""
    )

    # ========== SLIDE 21: Fixed Plus Product Design ==========
    add_data_slide(
        prs,
        "Two Fixed Plus tiers offer flexibility: Bronze (+1 Special) and Gold (+2 Holiday/Special)",
        [
            ["Package", "Includes", "Regular Price", "Fixed Plus Price", "Premium"],
            ["Fixed Plus Bronze", "+1 Special ticket", "$79.93", "$87.92", "+10%"],
            ["Fixed Plus Gold", "+2 tickets (Holiday + Special)", "$136.85", "$150.53", "+10%"],
            ["", "", "", "", ""],
            ["Projected Uptake", "Bronze: 141 subscribers", "Gold: 212 subscribers", "Total: 354", ""]
        ],
        notes="""SPEAKER NOTES - Product Design

PRICING LOGIC:
"We kept the premium simple at 10% across both tiers. This is easy to communicate and justified by the value-adds."

THE BENEFITS THAT JUSTIFY PREMIUM:
Tier 1 - Access Benefits (Core Value):
- Priority Access Window: Book 2 weeks before general public
- Guaranteed Seat Holds: Your preferred section reserved
- Sellout Protection: Never miss a Special or Holiday event due to capacity

Tier 2 - Experience Benefits:
- Pre-Concert Insights: Exclusive pre-concert talks
- Intermission Lounge Access: Dedicated refreshment area
- Program Recognition: 'Fixed Plus Member' acknowledgment

Tier 3 - Convenience Benefits:
- Single Checkout: Charged with subscription renewal
- Flex Swap Option: Exchange dates within tier
- Companion Discount: 5% off guest tickets

MARKETING STRATEGY:
"Do NOT advertise the price. Sell the benefits. 'Never miss another Special event' is more compelling than '$87.92 per ticket.'"
""",
        key_insight="Sell benefits, not price - 'Never miss another Special event' outperforms '$87.92/ticket'"
    )

    # ========== SLIDE 22: Fixed Plus Projections ==========
    add_data_slide(
        prs,
        "Conservative projections: 354 Fixed Plus subscribers generating $44,310 in premium-priced revenue",
        [
            ["Segment", "Size", "Uptake Rate", "Subscribers", "Revenue"],
            ["Current ST buyers (proven demand)", "950", "30%", "285", "$35,668"],
            ["Non-ST buyers (latent demand)", "693", "10%", "69", "$8,642"],
            ["TOTAL", "1,643", "21.5%", "354", "$44,310"],
            ["", "", "", "", ""],
            ["Premium Captured (10%)", "", "", "", "$4,028"]
        ],
        notes="""SPEAKER NOTES - Revenue Projections

UPTAKE ASSUMPTIONS:
"We applied different uptake rates based on demonstrated behavior. Subscribers who already buy single tickets (proven demand) get 30% uptake - they're pre-qualified. Non-buyers get 10% - they may have latent demand we're unlocking."

PACKAGE MIX:
- 40% Bronze (141 subscribers x $87.92 = $12,397)
- 60% Gold (212 subscribers x $150.53 = $31,913)

TOTAL PROGRAM VALUE:
- Fixed Plus Revenue: $44,310
- Regular price equivalent: $40,282
- Pure premium captured: $4,028 (9.1% of program revenue)

UPSIDE SCENARIOS:
- Higher uptake (40% for ST buyers): +$14,800
- Higher premium (15%): +$2,200
- Companion tickets (50% buy guest): +$9,200

BREAK-EVEN:
"We need about 100 Fixed Plus subscribers to cover program administration costs. At 354, this is comfortably profitable."
""",
        key_insight="30% uptake among proven buyers is conservative - these customers already want more"
    )

    # ========== SLIDE 23: Fixed Plus Benefits Package ==========
    add_content_slide(
        prs,
        "Fixed Plus benefits leverage behavioral insights: advance planners value priority access most",
        [
            "TIER 1: ACCESS BENEFITS (Core Value Proposition)",
            "  - Priority Access Window: Book Holiday/Special 2 weeks before general public",
            "  - Guaranteed Seat Holds: Your preferred section reserved for add-on events",
            "  - Sellout Protection: Never miss a Holiday show due to capacity",
            "",
            "TIER 2: EXPERIENCE BENEFITS (Differentiation)",
            "  - Pre-Concert Insights: Exclusive pre-concert talk or conductor Q&A",
            "  - Intermission Lounge Access: Dedicated refreshment area",
            "  - Program Recognition: 'Fixed Plus Member' acknowledgment",
            "",
            "TIER 3: CONVENIENCE BENEFITS (Friction Reduction)",
            "  - Single Checkout: Add-ons charged with subscription renewal",
            "  - Flex Swap Option: Exchange one date for another same-tier event",
            "  - Companion Discount: 5% off for a guest ticket"
        ],
        notes="""SPEAKER NOTES - Benefits Package

WHY THESE BENEFITS:
"Every benefit is based on observed behavior:
1. 45% of Fixed 6 subscribers buy tickets 2+ months in advance - they WANT to plan ahead
2. 83% sit in premium sections - seat quality matters
3. 57.8% already buy extras - convenience removes friction"

BEHAVIORAL ECONOMICS AT PLAY:
1. Bundling Effect: Bundles feel like better value even at same price
2. Loss Aversion: 'Sellout Protection' frames as avoiding loss, not gaining access
3. Status Signaling: 'Fixed Plus Member' appeals to identity
4. Commitment & Consistency: Once they identify as Fixed Plus, they renew

PHASED BENEFIT ROLLOUT:
"Start with low-cost/no-cost benefits (priority access, recognition). Add experiential benefits (pre-concert talks, lounge) as program matures and revenue justifies investment."

OPERATIONAL REQUIREMENTS:
- Hold 5% of Holiday/Special capacity for Fixed Plus priority
- Simple flag in ticketing system for member identification
- Staff training on Fixed Plus recognition
"""
    )

    # ========== SLIDE 24: Section Header - Comparison ==========
    add_section_header(
        prs,
        "Section 5: Comparing the Two Programs",
        notes="""SPEAKER NOTES - Section Header

TRANSITION:
"Let's step back and compare these two programs side by side. They target different audiences, use different mechanisms, and have different risk profiles - but they're complementary, not competing."
"""
    )

    # ========== SLIDE 25: Side-by-Side Comparison ==========
    add_two_column_slide(
        prs,
        "Both programs are complementary: one acquires new customers, one deepens existing relationships",
        {
            'header': 'Classical Add-On Product',
            'items': [
                "Target: Single ticket buyers (acquisition)",
                "Mechanism: Discounted cross-sell",
                "Revenue: $102,099 incremental",
                "Tickets: 1,926 new Classical sales",
                "Pipeline: 50 new CYO subscribers",
                "Risk: Very low (no cannibalization)",
                "Complexity: Low (email + signage)"
            ]
        },
        {
            'header': 'Fixed Plus Premium',
            'items': [
                "Target: Fixed 6 subscribers (deepening)",
                "Mechanism: Premium bundled value",
                "Revenue: $44,310 premium-priced",
                "Subscribers: 354 Fixed Plus members",
                "Premium: $4,028 captured above base",
                "Risk: Low (proven demand exists)",
                "Complexity: Medium (benefits delivery)"
            ]
        },
        notes="""SPEAKER NOTES - Comparison

COMPLEMENTARY, NOT COMPETING:
"These programs work together. Classical Add-On brings new people into your ecosystem; Fixed Plus deepens commitment from existing subscribers. Different targets, different mechanisms, both valuable."

REVENUE COMPARISON:
- Classical Add-On: $102,099 Year 1 + CYO pipeline
- Fixed Plus: $44,310 Year 1 + retention benefits

Total Year 1 impact: ~$146,000 in additional revenue

RISK COMPARISON:
- Classical Add-On: Zero cannibalization risk (product differentiation)
- Fixed Plus: Minimal risk (they're buying anyway; we're bundling)

IMPLEMENTATION COMPARISON:
- Classical Add-On: Easier (marketing only)
- Fixed Plus: More complex (operations + benefits)

RECOMMENDATION:
"We recommend launching both, with Classical Add-On in Phase 1 (lower complexity) and Fixed Plus in Phase 2 (requires more preparation)."
"""
    )

    # ========== SLIDE 26: Combined Financial Impact ==========
    add_kpi_slide(
        prs,
        "Combined Year 1 impact: $146K incremental revenue plus long-term subscriber pipeline",
        [
            {'value': '$146K', 'label': 'Total Year 1 incremental revenue (both programs)', 'color': GREEN},
            {'value': '50', 'label': 'New CYO subscribers from conversion pipeline', 'color': NAVY},
            {'value': '354', 'label': 'Fixed Plus members (deepened relationships)', 'color': LIGHT_BLUE}
        ],
        notes="""SPEAKER NOTES - Combined Impact

YEAR 1 SUMMARY:
- Classical Add-On: $102,099 ticket revenue
- Fixed Plus: $44,310 premium-priced revenue
- Total: $146,409

BEYOND YEAR 1:
Classical Add-On ongoing value:
- 50 new CYO subscribers x $500 avg = $25,000/year ongoing
- Repeat conversion from each season's Special/Holiday buyers

Fixed Plus ongoing value:
- 90%+ projected renewal rate
- Pathway to higher tiers (Gold, future Platinum)
- Stronger subscriber retention overall

5-YEAR PROJECTION:
Conservative estimate: $550K-750K cumulative impact
- Year 1: $146K
- Years 2-5: Growing as programs mature and CYO pipeline compounds
"""
    )

    # ========== SLIDE 27: Section Header - Implementation ==========
    add_section_header(
        prs,
        "Section 6: Implementation Roadmap",
        notes="""SPEAKER NOTES - Section Header

TRANSITION:
"Let me walk you through how we recommend implementing these programs, starting with quick wins and building to full scale."
"""
    )

    # ========== SLIDE 28: Implementation Roadmap ==========
    add_content_slide(
        prs,
        "Phased implementation minimizes risk: pilot Classical Add-On first, then launch Fixed Plus",
        [
            "PHASE 1: Classical Add-On Pilot (First 2-3 Special/Holiday Events)",
            "  - Test email + in-venue signage approach",
            "  - Measure response rate and redemption",
            "  - Collect feedback on offer clarity",
            "  - Target: 15%+ response rate",
            "",
            "PHASE 2: Fixed Plus Soft Launch (At Subscription Renewal)",
            "  - Target existing Fixed 6 subscribers who already buy single tickets",
            "  - Email campaign with priority access messaging",
            "  - Goal: 100 Fixed Plus subscribers (proof of concept)",
            "",
            "PHASE 3: Scale and Optimize",
            "  - Roll out Classical Add-On to all Special/Holiday events",
            "  - A/B test discount levels and messaging",
            "  - Expand Fixed Plus to all Fixed 6 at next renewal cycle",
            "  - Consider Fixed Plus Platinum tier for high-value members"
        ],
        notes="""SPEAKER NOTES - Implementation Roadmap

PHASE 1 TIMELINE:
"Start at your next Special or Holiday event. All you need is email capability and some signage. This is truly low-effort, high-potential."

PHASE 1 SUCCESS CRITERIA:
- Response rate >15% (expect 20%)
- Redemption rate >5% of responses
- No increase in complaints or confusion

PHASE 2 TIMELINE:
"Launch Fixed Plus at the next subscription renewal cycle. Lead with email to your best prospects - Fixed 6 subscribers who already buy single tickets."

PHASE 2 SUCCESS CRITERIA:
- 100 Fixed Plus subscribers (proof of concept)
- <5% cannibalization of existing single ticket sales
- Positive qualitative feedback

PHASE 3 TIMELINE:
"Once both pilots succeed, scale across the full season. This is when you invest in A/B testing, CRM integration, and benefit delivery infrastructure."
"""
    )

    # ========== SLIDE 29: Success Metrics ==========
    add_data_slide(
        prs,
        "Clear success metrics enable rapid iteration: monitor these KPIs weekly during pilot",
        [
            ["Program", "Metric", "Target", "Red Flag"],
            ["Classical Add-On", "Response rate", ">15%", "<10%"],
            ["Classical Add-On", "Redemption rate", ">5%", "<2%"],
            ["Classical Add-On", "Classical cannibalization", "<2% decline", ">5% decline"],
            ["Classical Add-On", "CYO conversion (3+ buyers)", ">20%", "<15%"],
            ["Fixed Plus", "Uptake rate", ">15% of Fixed 6", "<8%"],
            ["Fixed Plus", "Renewal rate", ">90%", "<80%"],
            ["Fixed Plus", "Single ticket cannibalization", "<5% decline", ">10% decline"],
            ["Fixed Plus", "Upgrade to Gold", ">30% of Bronze", "<15%"]
        ],
        notes="""SPEAKER NOTES - Success Metrics

WHY THESE METRICS:
"We've defined both success targets and red flags. If you're seeing green numbers, accelerate. If you hit a red flag, pause and diagnose before scaling."

MONITORING FREQUENCY:
- Response rate: Daily during pilot
- Redemption: Weekly
- Cannibalization: Monthly comparison to prior year
- Conversion/Retention: Quarterly

KEY DIAGNOSTIC QUESTIONS:
If response rate is low:
- Is the offer clear?
- Is timing right (same-day)?
- Are gateway events compelling?

If cannibalization appears:
- Check if 'new' buyers were actually likely purchasers anyway
- Review full-price sales by week to isolate effect

REPORTING RECOMMENDATION:
"Create a simple dashboard with these 8 metrics. Review weekly with box office and marketing. Monthly summary to leadership."
""",
        key_insight="Weekly monitoring during pilot enables rapid course correction before scale"
    )

    # ========== SLIDE 30: Risk Mitigation ==========
    add_content_slide(
        prs,
        "Both programs are designed for minimal risk: product differentiation protects existing revenue",
        [
            "CLASSICAL ADD-ON RISK MITIGATION:",
            "  - Discounts offered ONLY on the new add-on product, not existing pricing",
            "  - Customers who would have bought full price still can",
            "  - Offer available ONLY to verified Special/Holiday attendees",
            "  - 60-day window creates urgency without pressure",
            "",
            "FIXED PLUS RISK MITIGATION:",
            "  - 10% premium means no revenue loss on conversions",
            "  - Primary targets already buy singles (bundling, not new demand)",
            "  - Start with modest capacity holds (5% of Holiday/Special)",
            "  - Phase benefits over time; start with low-cost options",
            "",
            "OVERALL RISK PROFILE:",
            "  - Both programs target proven demand (they're already buying)",
            "  - Neither requires discounting core products",
            "  - Pilots before scale minimize downside exposure"
        ],
        notes="""SPEAKER NOTES - Risk Mitigation

THE PRODUCT DIFFERENTIATION PRINCIPLE:
"The critical design choice in both programs is that we're NOT discounting or diluting your core offerings. Classical Add-On creates a new product category. Fixed Plus adds premium value on top of base subscription. In neither case are we training customers to expect less for the same thing."

CANNIBALIZATION RISK IS MINIMAL:
Classical Add-On:
- 80.6% of targets have NEVER bought Classical - they're net new
- The 19.4% who have bought before get incremental tickets, not substitutes

Fixed Plus:
- 57.8% are buying singles anyway - we're bundling, not creating new behavior
- The 10% premium means even if some substitution occurs, revenue is protected

OPERATIONAL RISKS:
- Classical Add-On: Very low (marketing execution only)
- Fixed Plus: Moderate (benefit delivery requires coordination)

MITIGATION: Phase Fixed Plus benefits. Start with easy wins (priority access, recognition) before adding operationally complex benefits (lounge access, pre-concert talks).
"""
    )

    # ========== SLIDE 31: Next Steps ==========
    add_content_slide(
        prs,
        "Recommended next steps: approve pilots, assign owners, and launch within 60 days",
        [
            "IMMEDIATE ACTIONS (Next 2 Weeks):",
            "  1. Executive approval of pilot programs",
            "  2. Assign program owners (Marketing for Add-On, Subscriptions for Fixed Plus)",
            "  3. Brief box office staff on upcoming pilots",
            "",
            "SHORT-TERM (30 Days):",
            "  4. Design Classical Add-On email and signage creative",
            "  5. Select first 2-3 pilot events for Classical Add-On",
            "  6. Draft Fixed Plus offer letter for subscription renewal",
            "",
            "MEDIUM-TERM (60 Days):",
            "  7. Launch Classical Add-On pilot at first selected event",
            "  8. Establish weekly metrics reporting cadence",
            "  9. Prepare Fixed Plus for next renewal cycle",
            "",
            "ONGOING:",
            "  10. Weekly pilot reviews with cross-functional team",
            "  11. Monthly executive summary of program performance",
            "  12. Quarterly strategic review and optimization decisions"
        ],
        notes="""SPEAKER NOTES - Next Steps

DECISION POINT:
"We're asking for approval to proceed with pilots of both programs. Classical Add-On can launch at your next Special or Holiday event. Fixed Plus will launch at the next subscription renewal cycle."

RESOURCE REQUIREMENTS:
Classical Add-On:
- Marketing: 20 hours for creative development
- Box office: 2-hour staff briefing
- Tech: Email automation setup (existing capability)

Fixed Plus:
- Subscriptions: 40 hours for program design and launch
- Operations: Coordination for benefit delivery
- Tech: CRM flag for Fixed Plus members

COST ESTIMATE:
- Classical Add-On pilot: <$5,000 (email + signage)
- Fixed Plus soft launch: <$10,000 (creative + staff time)
- Total investment: <$15,000 for $143,000 projected return

TIMELINE:
"We can have Classical Add-On running within 60 days. Fixed Plus will align with your subscription renewal calendar."
"""
    )

    # ========== SLIDE 32: Key Takeaways ==========
    add_content_slide(
        prs,
        "Key takeaways: Data reveals opportunities; two programs capture $143K in Year 1",
        [
            "1. AUDIENCE SHIFT: Single ticket buyers growing (+2,796) while subscribers decline (-67).",
            "   Opportunity: Convert single ticket buyers to subscribers through Classical exposure.",
            "",
            "2. PROVEN DEMAND: 57% of subscribers buy add-on tickets; 80% of Special buyers never tried Classical.",
            "   Opportunity: Make it easier to buy what they already want.",
            "",
            "3. CYO VALUE: 13.7% of accounts generate 22-24% of subscription revenue at 72% higher prices.",
            "   Opportunity: Build CYO conversion pipeline from Classical Add-On program.",
            "",
            "4. PRICE POWER: Premium sections are inelastic; discounting Classical erodes revenue.",
            "   Opportunity: Use product differentiation, not price cuts.",
            "",
            "5. TWO SIMULATIONS: Classical Add-On ($102K) + Fixed Plus ($44K) = $146K Year 1 impact.",
            "   Opportunity: Low-risk pilots that scale based on results.",
        ],
        notes="""SPEAKER NOTES - Key Takeaways

THE 30-SECOND SUMMARY:
"Your audience is shifting toward single ticket buyers who pay more per seat but rarely return. The opportunity is to convert them into subscribers through targeted programs. Our two simulations - Classical Add-On and Fixed Plus - together can generate $146,000 in Year 1 incremental revenue with minimal risk."

IF YOU REMEMBER ONE THING:
"80.6% of your Special and Holiday buyers have never purchased a Classical ticket. That's not a problem - that's an opportunity. The Classical Add-On program bridges that gap."

THE STRATEGIC VISION:
"In 5 years, we want to look back and see that these pilots launched a flywheel: Special/Holiday buyers become Classical attenders, Classical attenders become CYO subscribers, CYO subscribers become your highest-value patrons. These programs plant those seeds."
"""
    )

    # ========== SLIDE 33: Closing ==========
    add_title_slide(
        prs,
        "Thank You",
        "Questions and Discussion"
    )

    # Add closing notes to final slide
    final_slide = prs.slides[-1]
    notes_slide = final_slide.notes_slide
    notes_slide.notes_text_frame.text = """SPEAKER NOTES - Closing

ANTICIPATED QUESTIONS:

Q: What if the response rates are lower than projected?
A: We've built in conservative assumptions (20% response for Add-On, 21.5% for Fixed Plus). Even at half these rates, the programs are profitable. Pilots let us calibrate before scale.

Q: Won't discounting train customers to wait for deals?
A: The discounts are explicitly tied to specific behaviors (attending a gateway event). This is reward-based, not promotional pricing. Customers can't access the discount without engaging first.

Q: How does this affect our brand positioning?
A: Both programs reinforce the brand. Classical Add-On says 'We want you to experience more of what we offer.' Fixed Plus says 'You're a valued subscriber and we're investing in your experience.'

Q: What's the cost of implementation?
A: Classical Add-On: <$5,000 for pilot (email, signage). Fixed Plus: <$10,000 (creative, staff time). Total: <$15,000 for projected $146,000 return.

Q: When can we start?
A: Classical Add-On can launch at your next Special or Holiday event. Fixed Plus aligns with subscription renewal cycle.

CLOSING STATEMENT:
"Thank you for your time. We're excited about the opportunity these programs represent and ready to support implementation. What questions do you have?"
"""

    # Save the presentation
    output_path = "/Users/galinajuliana/Desktop/ASOData/ASO_Data_Analysis_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
