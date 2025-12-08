#!/usr/bin/env python3
"""
Create PowerPoint presentation for Fixed Plus Premium Program Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

# Colors - using a premium gold/navy theme
NAVY = RGBColor(26, 42, 74)  # #1A2A4A
GOLD = RGBColor(212, 175, 55)  # #D4AF37
DARK_GRAY = RGBColor(68, 68, 68)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
LIGHT_GOLD = RGBColor(248, 241, 220)


def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text


def add_title_slide(prs, title, subtitle="", notes=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_section_slide(prs, title, notes=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_content_slide(prs, title, bullet_points, highlight_box=None, notes=""):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    # Optional highlight box
    if highlight_box:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(5.2), Inches(9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = GOLD
        box.line.fill.background()

        text_frame = box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = highlight_box
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_table_slide(prs, title, headers, rows, subtitle=None, notes=""):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Subtitle
    start_y = 1.0
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.font.italic = True
        start_y = 1.4

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_width = Inches(9)
    table_height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols,
                                    Inches(0.5), Inches(start_y),
                                    table_width, table_height).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

            # Highlight total row
            if "TOTAL" in str(row_data[0]).upper():
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD
                p.font.color.rgb = NAVY

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_metrics_slide(prs, title, metrics, notes=""):
    """Add a slide with key metrics in boxes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Metrics boxes
    num_metrics = len(metrics)
    box_width = 2.8
    box_height = 1.8
    spacing = 0.3
    start_x = (10 - (num_metrics * box_width + (num_metrics - 1) * spacing)) / 2

    for i, (label, value) in enumerate(metrics):
        x = start_x + i * (box_width + spacing)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(2), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()

        # Value
        value_box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(box_width), Inches(1))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(box_width), Inches(0.6))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_benefits_slide(prs, title, benefits_category, benefits, notes=""):
    """Add a slide showcasing benefits."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Category label
    cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
    tf = cat_box.text_frame
    p = tf.paragraphs[0]
    p.text = benefits_category
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.font.bold = True

    # Benefits
    start_y = 1.4
    for i, (benefit, description) in enumerate(benefits):
        y = start_y + i * 1.3

        # Benefit box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(y), Inches(9), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = GOLD

        # Benefit name
        name_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.1), Inches(8.5), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = benefit
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY

        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.5), Inches(8.5), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY

    if notes:
        add_slide_notes(slide, notes)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                    "Fixed Plus",
                    "A Premium Program for Your Most Loyal Subscribers",
                    notes="""OPENING TALKING POINTS:

Welcome everyone. Today I'm presenting Fixed Plus - a new premium program designed specifically for our most loyal subscribers.

The key insight driving this proposal is simple: our Fixed 6 subscribers want MORE from us, and they're already proving it with their wallets. We're just not making it easy for them.

Fixed Plus is about meeting that demand with a premium offering that delivers real value - and captures revenue we're currently leaving on the table.

Let me walk you through the data that led us here, the product design, and the projected impact.""")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Target: Fixed 6 subscribers - ASO's core subscription base (1,643 accounts)",
        "Insight: 57.8% already purchase additional single tickets (proven demand)",
        "Product: Premium bundle with Holiday/Special event access at 10% premium",
        "Strategy: Sell BENEFITS, not price - advance access, guaranteed seating, VIP status",
        "Revenue: $44,310 projected from 354 Fixed Plus subscribers"
    ], "Fixed 6 subscribers want more - Fixed Plus delivers",
    notes="""KEY POINTS TO EMPHASIZE:

1. TARGET MARKET: Fixed 6 is our largest subscription tier - 1,643 accounts representing 70% of all Fixed subscribers. These are committed patrons, not casual attendees.

2. THE INSIGHT: Here's what's remarkable - nearly 58% of these subscribers ALREADY go out of their way to buy additional single tickets. They're telling us with their behavior: "We want more than 6 concerts."

3. THE PRODUCT: Fixed Plus bundles what they already want - Holiday and Special event access - at a 10% premium. But we don't lead with price...

4. THE STRATEGY: We lead with BENEFITS. Priority access. Guaranteed seating. VIP status. The 10% premium feels justified when you're getting exclusive perks.

5. THE OPPORTUNITY: Conservative projections show $44,310 from just 354 subscribers - and this is a foundation we can build on.""")

    # Slide 3: Key Metrics - The Opportunity
    add_metrics_slide(prs, "The Opportunity: Fixed 6 Subscribers", [
        ("Total Fixed 6\nAccounts", "1,643"),
        ("Already Buy\nSingle Tickets", "57.8%"),
        ("Avg Extra\nSpend/Buyer", "$355")
    ],
    notes="""TALKING THROUGH THE NUMBERS:

Let me put these three numbers in context:

1,643 ACCOUNTS: This is a substantial base. Fixed 6 is our "gateway" subscription - serious enough to commit to 6 concerts, but not the heavy users who buy 18 or 24. They're the middle of the funnel.

57.8%: This is the critical insight. More than half of these subscribers are already buying additional single tickets on top of their subscription. They've demonstrated willingness to spend more with us. The question is: are we making it easy for them?

$355 AVERAGE: Among those who buy extra, they're spending $355 on average. That's significant incremental revenue per customer. And they're doing this despite friction - having to monitor ticket sales, compete for seats, go through separate checkout. Imagine if we removed that friction.

TRANSITION: Let's look at exactly what they're buying...""")

    # Slide 4: Behavior Evidence
    add_table_slide(prs, "Fixed 6 Single Ticket Behavior",
                    ["Event Type", "Unique Buyers", "Total Seats", "Revenue"],
                    [
                        ["Classical", "709", "3,274", "$208,590"],
                        ["Special", "381", "951", "$82,806"],
                        ["Holiday", "196", "708", "$41,980"]
                    ],
                    "381 Fixed 6 accounts bought Special tickets, 196 bought Holiday - demand exists",
                    notes="""ANALYZING THE DATA:

This table shows what Fixed 6 subscribers are buying as single tickets:

CLASSICAL ($208K): The largest category - 709 unique buyers purchasing Classical single tickets. This makes sense: their subscription is Classical-focused, and they want MORE Classical. This validates demand but isn't our Fixed Plus target.

SPECIAL ($82K): Here's where it gets interesting. 381 accounts - nearly a quarter of all Fixed 6 subscribers - bought Special event single tickets. These are premium experiences they're seeking out.

HOLIDAY ($42K): 196 accounts bought Holiday tickets. These events often sell out, creating urgency and FOMO.

THE INSIGHT: Combined, over 500 Fixed 6 accounts are buying Holiday and Special single tickets. That's the exact demand Fixed Plus is designed to capture - but with added benefits and at a premium price.

KEY QUESTION TO POSE: Why are we making our best customers compete with the general public for these seats?""")

    # Slide 5: What They Want (Evidence)
    add_content_slide(prs, "What Fixed 6 Subscribers Want", [
        "MORE EVENTS: 57.8% buy additional single tickets (avg 2.5 events)",
        "PLAN AHEAD: 45% purchase 2+ months in advance - advance access is valuable",
        "PREMIUM SEATING: 83% of subscriptions are in premium sections",
        "QUALITY: Average single ticket spend is $355 - willing to pay for value",
        "CONVENIENCE: They're already doing extra work to buy more tickets"
    ], "They want more - but current offerings don't serve them",
    notes="""BUILDING THE CASE - WHAT THE DATA TELLS US:

Let me translate these behaviors into customer needs:

1. MORE EVENTS: They're not satisfied with 6 concerts. The average extra buyer attends 2.5 additional events. They want a deeper relationship with ASO.

2. PLAN AHEAD: This is crucial for Fixed Plus design. 45% of their single ticket purchases happen 2+ months before the event. These are planners, not impulse buyers. "Priority Access" isn't just marketing fluff - it's something they genuinely value.

3. PREMIUM SEATING: 83% of their subscriptions are in Orchestra or Loge sections. They're not bargain hunters. They choose quality and are willing to pay for it.

4. WILLING TO PAY: $355 average spend on extras. These aren't price-sensitive customers looking for the cheapest option.

5. CONVENIENCE: Think about what they're doing now - monitoring on-sale dates, competing for seats, going through separate purchases. They're doing this despite friction. What if we removed that friction and added VIP perks?

THE BOTTOM LINE: Fixed 6 subscribers are telling us what they want through their behavior. Fixed Plus is simply packaging it better.""")

    # Slide 6: Section - The Solution
    add_section_slide(prs, "The Solution: Fixed Plus",
                      notes="""TRANSITION:

Now that we've established the demand exists - and it clearly does - let me show you how Fixed Plus is designed to meet that demand while capturing premium revenue.

The key principle here is PRODUCT DIFFERENTIATION. We're not discounting. We're not cannibalizing. We're creating a new premium tier that delivers genuine value.""")

    # Slide 7: Product Design
    add_table_slide(prs, "Fixed Plus Package Options",
                    ["Package", "Includes", "Premium"],
                    [
                        ["Fixed Plus Bronze", "+1 Special ticket", "10%"],
                        ["Fixed Plus Gold", "+2 tickets (Holiday + Special)", "10%"]
                    ],
                    "Bundles what they already want - at their subscription renewal",
                    notes="""PRODUCT STRUCTURE:

We're offering two tiers to capture different levels of demand:

FIXED PLUS BRONZE:
- Adds 1 Special ticket to their subscription
- 10% premium over regular single ticket price
- Entry-level option for those who want to "try" the program
- Price: approximately $87.92 (vs $79.93 regular)

FIXED PLUS GOLD:
- Adds 2 tickets: 1 Holiday + 1 Special event
- 10% premium on both
- For the more engaged subscriber who wants the full experience
- Price: approximately $150.53 total

WHY 10% PREMIUM?
- High enough to capture meaningful revenue
- Low enough to feel justified by the benefits
- Positions as "premium" without feeling exploitative

TIMING IS KEY: We offer this at subscription renewal - when they're already in "buying mode" and thinking about their ASO commitment for the year. It's a natural upsell moment.""")

    # Slide 8: Benefits Tier 1
    add_benefits_slide(prs, "Fixed Plus Benefits", "ACCESS BENEFITS (Core Value)", [
        ("Priority Access Window", "Book Holiday/Special events 2 weeks before general public"),
        ("Guaranteed Seat Holds", "Your preferred section reserved for add-on events"),
        ("Sellout Protection", "Never miss a Holiday show due to capacity")
    ],
    notes="""ACCESS BENEFITS - THE CORE VALUE PROPOSITION:

These three benefits are the foundation of Fixed Plus. They solve real problems our subscribers face:

PRIORITY ACCESS WINDOW (2 weeks early):
- Remember: 45% of Fixed 6 buyers purchase 2+ months in advance
- They're planners who want certainty
- Two weeks early access means they get first pick of dates and seats
- Creates genuine exclusivity

GUARANTEED SEAT HOLDS:
- 83% sit in premium sections for their subscription
- When buying single tickets, they risk not getting comparable seats
- Fixed Plus guarantees their preferred section is available
- Removes the anxiety of "will I get good seats?"

SELLOUT PROTECTION:
- Holiday events often sell out
- Nothing is more frustrating than wanting to attend and being shut out
- This benefit uses LOSS AVERSION psychology - avoiding the pain of missing out
- Frame it as: "Never miss a Holiday concert again"

These aren't invented benefits - they directly address friction points we see in the data.""")

    # Slide 9: Benefits Tier 2
    add_benefits_slide(prs, "Fixed Plus Benefits", "EXPERIENCE BENEFITS (Differentiation)", [
        ("Pre-Concert Insights", "Exclusive pre-concert talk or conductor Q&A"),
        ("Intermission Lounge Access", "Dedicated refreshment area - skip the lines"),
        ("Program Recognition", "'Fixed Plus Member' acknowledgment - VIP status")
    ],
    notes="""EXPERIENCE BENEFITS - CREATING DIFFERENTIATION:

These benefits transform Fixed Plus from a "ticket bundle" into a "membership experience":

PRE-CONCERT INSIGHTS:
- Exclusive educational content before Holiday/Special events
- Could be a conductor Q&A, musicologist talk, or behind-the-scenes preview
- Makes them feel like insiders, not just ticket holders
- Relatively low cost to deliver, high perceived value

INTERMISSION LOUNGE ACCESS:
- Dedicated refreshment area (if venue supports it)
- Avoids the crowded bar lines
- Creates a physical "members only" space
- Note: This may require operational setup - can phase in later

PROGRAM RECOGNITION:
- "Fixed Plus Member" acknowledgment in programs or communications
- Appeals to identity and status
- Makes the membership visible and social
- Simple to implement, meaningful to recipients

PHASING: We can launch with access benefits first, then add experience benefits as the program matures. Start simple, prove the concept, then enhance.""")

    # Slide 10: Benefits Tier 3
    add_benefits_slide(prs, "Fixed Plus Benefits", "CONVENIENCE BENEFITS (Remove Friction)", [
        ("Single Checkout", "Add-on events charged with subscription renewal"),
        ("Flex Swap Option", "Exchange one add-on date for another same-tier event"),
        ("Companion Discount", "5% off for a guest ticket on add-on events")
    ],
    notes="""CONVENIENCE BENEFITS - REMOVING FRICTION:

These benefits make it EASY to say yes to Fixed Plus:

SINGLE CHECKOUT:
- Add-on events charged with subscription renewal
- No separate transactions to manage
- No need to remember on-sale dates
- One decision, one payment, done
- This alone removes significant friction

FLEX SWAP OPTION:
- Life happens - schedules change
- Allow exchange of one add-on date for another event in the same tier
- Reduces the risk of commitment
- "What if I can't make that date?" is no longer a barrier
- Builds trust and reduces buyer hesitation

COMPANION DISCOUNT:
- 5% off for a guest ticket on add-on events
- Data shows many buyers purchase 2+ seats
- Makes it easier to bring a friend or partner
- Small discount, but signals "we value you sharing this experience"

THE PSYCHOLOGY: Every friction point is an opportunity for the customer to say "maybe later" or "too complicated." These benefits systematically remove those friction points.""")

    # Slide 11: Marketing Strategy
    add_content_slide(prs, "Marketing: Sell Benefits, Not Price", [
        "DO NOT advertise the 10% premium",
        "INSTEAD, lead with exclusive benefits:",
        "    'Never miss a Holiday concert - guaranteed access'",
        "    'Your seats. Your events. Your way.'",
        "    'Join the Fixed Plus family'",
        "Position as an UPGRADE, not an expense",
        "Offer at subscription renewal - natural decision point"
    ],
    notes="""CRITICAL MARKETING GUIDANCE:

This is important - how we position Fixed Plus determines its success:

DO NOT LEAD WITH PRICE:
- Never say "for just 10% more..."
- Price-focused messaging attracts price-focused objections
- We want value-focused customers

LEAD WITH BENEFITS:
- "Never miss a Holiday concert" → Sellout protection
- "Your seats. Your events. Your way." → Control and choice
- "Join the Fixed Plus family" → Belonging and status

POSITIONING AS UPGRADE:
- It's not "paying extra" - it's "upgrading your experience"
- Like first class vs economy: same destination, better journey
- Appeals to aspiration, not calculation

TIMING - SUBSCRIPTION RENEWAL:
- This is when they're already engaged and deciding
- Natural moment to present options
- "Would you like to enhance your subscription this year?"
- Much easier than cold outreach

MESSAGING EXAMPLES:
- Email subject: "Exclusive: Fixed Plus membership now open"
- NOT: "Add Holiday tickets for 10% premium"

The 10% premium should feel like a natural part of a premium experience, not the defining feature.""")

    # Slide 12: Section - Projections
    add_section_slide(prs, "Revenue Projections",
                      notes="""TRANSITION:

Now let's look at the numbers. These projections are based on conservative assumptions about uptake, grounded in the behavioral data we've analyzed.

I'll show you our target segments, expected conversion rates, and the resulting revenue impact.""")

    # Slide 13: Target Segments
    add_table_slide(prs, "Target Segments & Uptake",
                    ["Segment", "Size", "Uptake Rate", "Converts"],
                    [
                        ["Current ST buyers (proven demand)", "950", "30%", "285"],
                        ["Non-ST buyers (latent demand)", "693", "10%", "69"],
                        ["TOTAL", "1,643", "-", "354"]
                    ],
                    notes="""UPTAKE ASSUMPTIONS EXPLAINED:

We're targeting two segments with different expected conversion rates:

SEGMENT 1: CURRENT SINGLE TICKET BUYERS (950 accounts)
- These are Fixed 6 subscribers who ALREADY buy extra tickets
- They've demonstrated the behavior we're trying to capture
- 30% uptake assumption is reasonable because:
  - They already want what we're offering
  - We're adding benefits (priority access, convenience)
  - We're removing friction (single checkout)
- 30% of 950 = 285 expected converts

SEGMENT 2: NON-SINGLE TICKET BUYERS (693 accounts)
- These Fixed 6 subscribers haven't bought extras YET
- Could be due to friction, awareness, or lower interest
- 10% uptake assumption is conservative because:
  - They haven't shown the behavior
  - But benefits might unlock latent demand
  - "Sellout protection" might appeal to FOMO
- 10% of 693 = 69 expected converts

TOTAL: 354 Fixed Plus subscribers

WHY THESE RATES ARE CONSERVATIVE:
- Industry benchmarks for subscription upsells: 15-40%
- Our benefits directly address documented pain points
- Upside potential if marketing resonates""")

    # Slide 14: Revenue Summary
    add_table_slide(prs, "Projected Revenue",
                    ["Package", "Subscribers", "Revenue"],
                    [
                        ["Fixed Plus Bronze (+1 Special)", "141", "$12,397"],
                        ["Fixed Plus Gold (+2 Holiday/Special)", "212", "$31,913"],
                        ["TOTAL", "354", "$44,310"]
                    ],
                    notes="""REVENUE BREAKDOWN:

PACKAGE MIX ASSUMPTION:
- 40% choose Bronze (entry-level, 1 Special ticket)
- 60% choose Gold (fuller experience, 2 tickets)
- This mix is based on typical tiered pricing behavior: most choose the middle/premium option when benefits are clear

BRONZE REVENUE: $12,397
- 141 subscribers × $87.92 (Special at 10% premium)
- Entry point for cautious buyers
- Many may upgrade to Gold in year 2

GOLD REVENUE: $31,913
- 212 subscribers × ~$150.53 (Holiday + Special at 10% premium)
- The "sweet spot" offering
- Higher value, higher engagement

TOTAL: $44,310
- First-year revenue from Fixed Plus program
- Remember: this is INCREMENTAL - either new revenue or premium capture
- Foundation for growth in subsequent years

IMPORTANT CONTEXT: This is conservative. If we achieve 40% uptake among ST buyers instead of 30%, that's an additional $14,800+ in revenue.""")

    # Slide 15: Financial Metrics
    add_metrics_slide(prs, "Financial Summary", [
        ("Total Fixed Plus\nRevenue", "$44,310"),
        ("Premium Captured\n(10%)", "$4,028"),
        ("Avg Revenue per\nFixed Plus Sub", "$125")
    ],
    notes="""FINANCIAL SUMMARY:

Three numbers to remember:

$44,310 TOTAL REVENUE:
- This is what Fixed Plus generates in year one
- For context: if these customers bought the same tickets at regular price, we'd get $40,282
- Fixed Plus adds both volume AND margin

$4,028 PREMIUM CAPTURED:
- This is the pure premium - the 10% above regular pricing
- This is revenue we would NOT have captured otherwise
- Relatively small in year one, but grows with program scale
- Also: doesn't account for customers who wouldn't have bought at all

$125 AVERAGE REVENUE PER SUBSCRIBER:
- Each Fixed Plus member generates $125 in add-on revenue
- Compare to: current average of $355 for those who buy ST (but only 57.8% do)
- Fixed Plus captures more people at a more predictable rate

THE BIGGER PICTURE:
- Year 1: $44,310 (354 subscribers)
- Year 2 with 90% renewal + new: potentially $65K+
- Year 3 with Platinum tier: $90K+
- This is a foundation for sustainable growth""")

    # Slide 16: Why This Works
    add_content_slide(prs, "Why This Works (Psychology)", [
        "BUNDLING: Bundles feel like better value than buying separately",
        "LOSS AVERSION: 'Sellout Protection' = avoiding loss (powerful motivator)",
        "STATUS: 'Fixed Plus Member' appeals to identity and belonging",
        "COMMITMENT: Once a 'Plus' member, more likely to renew & upgrade",
        "CONVENIENCE: Single checkout removes friction that blocks purchases"
    ],
    notes="""BEHAVIORAL ECONOMICS BEHIND FIXED PLUS:

Let me explain WHY these design choices work psychologically:

BUNDLING EFFECT:
- Research shows customers perceive bundles as better value
- Even at the same total price, "package" feels like a deal
- Fixed Plus positions tickets + benefits as integrated offering

LOSS AVERSION (most powerful):
- Psychologically, avoiding a loss is 2x more motivating than achieving a gain
- "Sellout Protection" = avoiding the PAIN of missing out
- "Priority Access" = avoiding the FRUSTRATION of bad seats
- We're not selling tickets; we're selling peace of mind

STATUS & IDENTITY:
- "Fixed Plus Member" creates an identity
- People act consistently with their self-image
- Once they see themselves as a "Plus member," they behave like one
- Recognition in programs reinforces this identity

COMMITMENT & CONSISTENCY:
- Once someone upgrades to Fixed Plus, they're invested
- They'll want to "get their money's worth" from benefits
- Renewal becomes natural - they've already self-identified as Plus members

CONVENIENCE PREMIUM:
- People pay more for convenience (Amazon Prime, priority boarding)
- Single checkout removes decision fatigue
- Friction is the enemy of conversion""")

    # Slide 17: Risk Mitigation
    add_content_slide(prs, "Risk Mitigation", [
        "CANNIBALIZATION: 10% premium means no revenue loss on conversions",
        "If they would have bought anyway → we capture extra margin",
        "If they wouldn't have bought → we unlock new revenue",
        "OPERATIONAL: Start with modest seat holds (5% of Holiday/Special)",
        "BENEFITS: Phase in over time - start with access, add experience later"
    ],
    notes="""ADDRESSING POTENTIAL CONCERNS:

CANNIBALIZATION RISK:
Q: "Won't this just shift revenue we would have gotten anyway?"
A: Two scenarios, both positive:

Scenario 1 - They would have bought single tickets anyway:
- Now they pay 10% premium
- We capture extra margin on inevitable purchases
- Net positive

Scenario 2 - They wouldn't have bought:
- Fixed Plus unlocks demand that friction was blocking
- 100% incremental revenue
- Net positive

There's no scenario where we lose.

OPERATIONAL CONCERNS:
Q: "How do we manage priority access inventory?"
A: Start conservatively:
- Hold 5% of Holiday/Special capacity for Fixed Plus
- Monitor demand and adjust
- If underutilized, release to general sale
- Simple ticketing system flag

BENEFIT DELIVERY:
Q: "Can we actually deliver on pre-concert talks, lounge access?"
A: Phase the benefits:
- Year 1: Focus on ACCESS benefits (easy to deliver)
- Year 2: Add EXPERIENCE benefits (as program proves itself)
- Year 3: Premium tier with exclusive experiences
- Don't over-promise; build credibility gradually""")

    # Slide 18: Success Metrics
    add_table_slide(prs, "Success Metrics",
                    ["Metric", "Target", "Red Flag"],
                    [
                        ["Fixed Plus uptake rate", ">15% of Fixed 6", "<8%"],
                        ["Renewal rate (Fixed Plus)", ">90%", "<80%"],
                        ["Upgrade Bronze → Gold", ">30%", "<15%"],
                        ["Net Promoter Score", ">60", "<40"]
                    ],
                    notes="""HOW WE'LL MEASURE SUCCESS:

UPTAKE RATE (>15% target):
- Our projection assumes ~22% overall uptake
- 15% is the minimum threshold for success
- Below 8% signals fundamental product-market fit issues
- Track by segment: ST buyers vs. non-ST buyers

RENEWAL RATE (>90% target):
- Critical for long-term value
- High renewal = satisfied customers
- Below 80% means benefits aren't delivering perceived value
- Survey churned members to understand why

UPGRADE BRONZE → GOLD (>30% target):
- Shows customers want MORE after trying Bronze
- Validates the tiered approach
- Below 15% might mean Gold benefits aren't compelling enough
- Opportunity to refine Gold offering

NET PROMOTER SCORE (>60 target):
- Would they recommend Fixed Plus to fellow subscribers?
- High NPS = organic growth through word-of-mouth
- Below 40 indicates satisfaction problems
- Use feedback to improve benefits

REVIEW CADENCE: Monthly tracking, quarterly deep-dive, annual strategic review""")

    # Slide 19: Implementation
    add_content_slide(prs, "Implementation Roadmap", [
        "PHASE 1 - Pilot:",
        "    Target existing ST buyers (950 accounts) at renewal",
        "    Focus on access benefits, goal: 100 subscribers",
        "PHASE 2 - Expand:",
        "    Target all Fixed 6, add experience benefits",
        "    Goal: 250 subscribers",
        "PHASE 3 - Premium Tier:",
        "    Introduce 'Fixed Plus Platinum' with backstage access"
    ],
    notes="""IMPLEMENTATION ROADMAP:

PHASE 1 - PILOT (First renewal cycle):
Target: 950 accounts who already buy single tickets
Why them first:
- Proven demand (they already do what we're offering)
- Lower risk of rejection
- Quick proof of concept

Approach:
- Personal email from subscriber services
- Focus on ACCESS benefits (priority, guaranteed seats, sellout protection)
- Single checkout convenience

Goal: 100 subscribers (10.5% of target segment)
Success criteria: >15% uptake, >85% satisfaction

PHASE 2 - EXPAND (Second renewal cycle):
Target: All 1,643 Fixed 6 accounts
Additions:
- Experience benefits (pre-concert talks, lounge access)
- Refined messaging based on Phase 1 learnings
- Testimonials from Phase 1 members

Goal: 250 total subscribers
Success criteria: >90% Phase 1 renewal, >12% new uptake

PHASE 3 - PREMIUM (Year 2-3):
Introduce "Fixed Plus Platinum":
- All Gold benefits PLUS
- Backstage access / meet the conductor
- Premium pricing (25% above Gold)
- Target: Top 20% of Gold members

This creates an upgrade path and maximizes lifetime value.""")

    # Slide 20: Conclusion
    add_content_slide(prs, "Conclusion", [
        "Fixed 6 subscribers are loyal but UNDERSERVED",
        "57.8% already buy extra tickets - demand is PROVEN",
        "Fixed Plus meets them with bundled convenience + VIP benefits",
        "10% premium is justified by exclusive access & experience",
        "Projected: 354 subscribers, $44,310 revenue, pathway to expand"
    ], "Give your best subscribers what they want - and capture the premium",
    notes="""CLOSING SUMMARY:

Let me leave you with five key takeaways:

1. THE OPPORTUNITY IS REAL
Fixed 6 subscribers are our most loyal yet underserved segment. They've told us through their behavior that they want more.

2. THE DEMAND IS PROVEN
57.8% already buy single tickets. We're not creating demand; we're packaging it better.

3. THE PRODUCT IS DESIGNED FOR THEM
Fixed Plus isn't a generic upsell. Every benefit addresses a specific pain point we identified in the data: advance planning, seat quality, convenience.

4. THE PREMIUM IS JUSTIFIED
10% is modest when you're getting priority access, guaranteed seating, sellout protection, and VIP status. We're not gouging; we're delivering value.

5. THE NUMBERS WORK
354 subscribers, $44,310 year-one revenue, with a clear pathway to grow through renewals, upgrades, and premium tiers.

THE ASK: Approve the pilot program targeting our 950 single-ticket-buying Fixed 6 subscribers at their next renewal cycle.

The worst case: We learn and adjust.
The best case: We build a sustainable premium revenue stream while deepening relationships with our most valuable patrons.

Questions?""")

    # Slide 21: Thank You
    add_title_slide(prs, "Thank You", "Questions?",
                    notes="""PREPARED ANSWERS FOR LIKELY QUESTIONS:

Q: Why 10% premium specifically?
A: Balances revenue capture with perceived value. High enough to be meaningful, low enough to feel justified by benefits. We can test higher premiums in Phase 3.

Q: What if Holiday/Special events don't sell out?
A: Fixed Plus still delivers value through priority access, guaranteed seating, and convenience. Sellout protection is insurance - valuable even if not needed.

Q: How do we handle existing subscribers who already renewed?
A: Offer mid-cycle upgrade option with prorated pricing. Don't force them to wait a full year.

Q: What's the cost to deliver these benefits?
A: Access benefits: Near-zero (system configuration). Experience benefits: Modest (speaker fees, lounge setup). ROI is strong given $40K+ revenue.

Q: Could this cannibalize our Holiday/Special single ticket sales?
A: Unlikely. Fixed Plus targets EXISTING subscribers, not the general public. General public still buys normally.

Q: What if uptake is lower than projected?
A: Phase 1 pilot is designed to test this. If <8% uptake, we reassess benefits and messaging before Phase 2.

THANK YOU FOR YOUR TIME. I'm happy to dive deeper into any aspect of this analysis.""")

    # Save
    output_path = Path(__file__).resolve().parent / "Fixed_Plus_Premium_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
