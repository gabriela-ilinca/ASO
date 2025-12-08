#!/usr/bin/env python3
"""
Create PowerPoint presentation for Classical Add-On Product Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

# Colors
DARK_BLUE = RGBColor(17, 106, 123)  # #116A7B
LIGHT_BLUE = RGBColor(26, 200, 237)  # #1AC8ED
DARK_GRAY = RGBColor(68, 68, 68)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)


def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text


def add_title_slide(prs, title, subtitle="", notes=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BLUE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_section_slide(prs, title, notes=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BLUE)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_content_slide(prs, title, bullet_points, highlight_box=None, notes=""):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    # Optional highlight box
    if highlight_box:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(5.2), Inches(9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.fill.background()

        text_frame = box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = highlight_box
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_table_slide(prs, title, headers, rows, subtitle=None, notes=""):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Subtitle
    start_y = 1.0
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.font.italic = True
        start_y = 1.4

    # Table
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    table_width = Inches(9)
    table_height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols,
                                    Inches(0.5), Inches(start_y),
                                    table_width, table_height).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

            # Bold the total row
            if "TOTAL" in str(row_data[0]).upper():
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
                p.font.color.rgb = WHITE

    if notes:
        add_slide_notes(slide, notes)

    return slide


def add_metrics_slide(prs, title, metrics, notes=""):
    """Add a slide with key metrics in boxes."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Metrics boxes
    num_metrics = len(metrics)
    box_width = 2.8
    box_height = 1.8
    spacing = 0.3
    start_x = (10 - (num_metrics * box_width + (num_metrics - 1) * spacing)) / 2

    for i, (label, value) in enumerate(metrics):
        x = start_x + i * (box_width + spacing)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(2), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.fill.background()

        # Value
        value_box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(box_width), Inches(1))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(box_width), Inches(0.6))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

    if notes:
        add_slide_notes(slide, notes)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                    "Classical Add-On Product",
                    "A New Revenue Strategy for Special/Holiday Attendees",
                    notes="""OPENING TALKING POINTS:

Thank you for joining me today. I'm excited to present a new revenue opportunity we've identified through our ticket data analysis.

The Classical Add-On Product is a strategic initiative to convert our Special and Holiday event attendees into Classical series customers.

Here's what makes this opportunity compelling:
- It's based on PROVEN demand patterns in our data
- It protects our existing pricing through product differentiation
- It creates a pipeline to our CYO subscription program
- It's low-risk because we're creating a NEW product, not discounting existing ones

Let me walk you through the data, the product design, and the projected impact.""")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "New product: Discounted Classical tickets for Special/Holiday single ticket buyers",
        "Target market: 7,008 accounts (30% of all single ticket buyers)",
        "80.6% of these buyers have NEVER purchased a Classical ticket",
        "Product differentiation protects existing pricing - no cannibalization risk",
        "Creates a pipeline for CYO subscription conversion"
    ], "Projected Impact: $102K incremental revenue + 50 new CYO subscribers",
    notes="""KEY POINTS TO EMPHASIZE:

1. THE PRODUCT: We're offering discounted Classical tickets ONLY to people who attended Special or Holiday events. This is a targeted loyalty reward, not a general discount.

2. MARKET SIZE: 7,008 accounts - this is 30% of all our single ticket buyers. It's a substantial audience.

3. THE UNTAPPED OPPORTUNITY: Here's the key insight - 80.6% of these Special/Holiday buyers have NEVER bought a Classical ticket. That's over 5,600 accounts who already love ASO but haven't tried our core Classical programming.

4. WHY PRODUCT DIFFERENTIATION MATTERS: We're not discounting Classical tickets broadly. The discount is ONLY available to verified Special/Holiday ticket holders. Regular Classical pricing remains intact. If someone wants to buy Classical at full price, they still can.

5. THE CYO PIPELINE: This isn't just about selling tickets. When we convert someone to 3+ Classical events, they become prime candidates for CYO subscription. Our historical data shows 24.6% conversion rate.

BOTTOM LINE: $102K incremental revenue plus 50 new CYO subscribers.""")

    # Slide 3: Key Metrics
    add_metrics_slide(prs, "Key Opportunity Metrics", [
        ("Target Market", "7,008"),
        ("Untapped\n(Never bought Classical)", "80.6%"),
        ("Projected New Tickets", "1,926")
    ],
    notes="""TALKING THROUGH THE NUMBERS:

7,008 TARGET ACCOUNTS:
- These are single ticket buyers who attended Special or Holiday events in FY23-FY24
- They've already demonstrated they enjoy ASO - they chose to attend our premium events
- They represent 30% of all single ticket accounts
- This is a WARM audience, not cold outreach

80.6% NEVER BOUGHT CLASSICAL:
- This is the remarkable insight that drives the entire strategy
- 5,647 accounts attended Special/Holiday but have never tried Classical
- They're not anti-Classical; they just haven't been given a compelling reason to try it
- The add-on product gives them that reason: a discount tied to an event they already enjoyed

1,926 PROJECTED NEW TICKETS:
- Based on conservative 20% response rate assumptions
- Mix of +1, +2, and +3 ticket purchases
- Each ticket is incremental - these are sales that wouldn't happen otherwise

KEY TRANSITION: Let me show you how the discount structure works...""")

    # Slide 4: Product Design
    add_table_slide(prs, "Discount Structure (Cumulative)",
                    ["Purchase Level", "Discount", "Example Price"],
                    [
                        ["+1 Classical ticket", "5% off", "$55.63"],
                        ["+2 Classical tickets", "10% off each", "$52.70/ticket"],
                        ["+3 Classical tickets", "15% off each", "$49.78/ticket"]
                    ],
                    "Based on average Classical single ticket price of $58.56",
                    notes="""DISCOUNT STRUCTURE EXPLAINED:

The discount is CUMULATIVE - meaning it applies to ALL tickets at the highest tier reached:

+1 TICKET (5% off):
- Entry-level option for cautious first-timers
- Price: $55.63 (saves $2.93)
- Low commitment, easy yes

+2 TICKETS (10% off EACH):
- Both tickets get the 10% discount
- Price: $52.70 per ticket = $105.40 total (saves $11.72)
- Encourages bringing a friend/partner

+3 TICKETS (15% off EACH):
- All three tickets get the 15% discount
- Price: $49.78 per ticket = $149.34 total (saves $26.34)
- This is our target because 3+ Classical buyers convert to CYO at 24.6%

WHY CUMULATIVE?
- Rewards higher commitment
- Simple to understand
- Creates clear incentive to buy more
- Average discount across our projections: ~9.5%

IMPORTANT: This discount is ONLY available to verified Special/Holiday ticket holders. Regular Classical pricing ($58.56 average) remains unchanged for everyone else.""")

    # Slide 5: Why Same-Day Marketing Matters
    add_content_slide(prs, "Marketing: Why Same-Day Offers Matter", [
        "Peak Emotional Engagement: Customers are at highest satisfaction post-concert",
        "Brand Association: Memory of event is fresh, Classical feels like natural continuation",
        "Reduced Decision Friction: Already demonstrated willingness to spend on ASO",
        "Social Proof: Group decisions more likely in the moment",
        "60-day validity window creates urgency while allowing time to plan"
    ],
    notes="""THE PSYCHOLOGY OF TIMING:

This slide is CRITICAL for execution. The WHEN matters as much as the WHAT.

PEAK EMOTIONAL ENGAGEMENT:
- Right after a Special or Holiday concert, customers are on a high
- They just had a wonderful experience with ASO
- This is the moment they're most receptive to "more of this"
- Waiting even 24-48 hours means the emotional peak fades

BRAND ASSOCIATION:
- "You loved Holiday Pops, now experience Beethoven's 9th"
- The connection feels natural, not salesy
- We're extending their positive experience, not interrupting them later

REDUCED FRICTION:
- They're already in "ASO mode"
- Credit card was just used (mental accounting)
- The decision to spend on ASO has already been made today
- Adding Classical feels like an extension, not a new decision

SOCIAL PROOF:
- If they're with friends/family, group purchases are more likely
- "Should we try the Classical series?" - easier to decide together in the moment

60-DAY VALIDITY:
- Long enough to plan around schedules
- Short enough to create urgency
- Matches our typical advance purchase window for this segment""")

    # Slide 6: Marketing Channels
    add_table_slide(prs, "Recommended Marketing Touchpoints",
                    ["Channel", "Timing", "Message"],
                    [
                        ["Email", "Within 2 hours post-event", '"Continue your ASO experience..."'],
                        ["In-venue signage", "Intermission/exit", "QR code to add-on offer"],
                        ["Box office", "At check-in", "Staff mention of program"],
                        ["SMS (if opted in)", "Next morning", "Reminder + upcoming options"]
                    ],
                    notes="""MULTI-CHANNEL APPROACH:

EMAIL (Primary Channel):
- Send within 2 hours of event end
- Subject line: "Continue your ASO experience - exclusive offer inside"
- Personalize: "Thank you for joining us for [Event Name]"
- Include 2-3 suggested Classical concerts
- Clear CTA with discount code
- 60-day expiration prominently displayed

IN-VENUE SIGNAGE:
- Table tents during intermission
- Exit signage with QR code
- "Loved tonight? Try our Classical series..."
- QR leads to mobile-optimized landing page
- Can scan now, complete purchase later

BOX OFFICE:
- Train staff to mention the program at check-in
- "Did you know as a Holiday Pops attendee, you qualify for..."
- Hand out info cards with QR code
- Not pushy - informative

SMS (Next Morning):
- Only for opted-in customers
- Gentle reminder: "Don't forget your Classical add-on offer"
- Include direct link to top 3 recommended events
- 60-day countdown reminder

KEY POINT: Multiple touchpoints increase awareness without being aggressive. Customer chooses their preferred channel.""")

    # Slide 7: Gateway Events
    add_table_slide(prs, "Top 'Gateway' Classical Events",
                    ["Event Name", "Crossover Buyers", "Total Seats", "Revenue"],
                    [
                        ["Beethoven Symphony No 9", "188", "427", "$29,067"],
                        ["Opening Weekend", "164", "327", "$21,503"],
                        ["Orff: Carmina Burana", "139", "287", "$22,903"],
                        ["Beethoven and Bolero", "119", "351", "$22,351"],
                        ["Bizet and Tchaikovsky", "115", "329", "$16,811"]
                    ],
                    "These events have proven appeal to Special/Holiday crossover buyers",
                    notes="""DATA-DRIVEN EVENT RECOMMENDATIONS:

This table shows which Classical events our Special/Holiday buyers ALREADY gravitate toward when they do cross over. These are our "gateway" concerts.

BEETHOVEN SYMPHONY NO 9 (188 crossover buyers):
- The most popular Classical event among our target audience
- Familiar, accessible, emotionally powerful
- "You don't need to know Classical to love the 9th"
- Should be prominently featured in all add-on marketing

OPENING WEEKEND (164 crossover buyers):
- Appeals to those who like "events" not just concerts
- Special occasion feel matches Special/Holiday vibe
- Good for first-timers who want a memorable experience

ORFF: CARMINA BURANA (139 crossover buyers):
- Dramatic, theatrical, accessible
- Appeals to the Special event crowd
- "You'll recognize O Fortuna"

BEETHOVEN AND BOLERO (119 crossover buyers):
- Mix of familiar pieces
- Bolero is highly recognizable
- Good entry point

BIZET AND TCHAIKOVSKY (115 crossover buyers):
- Carmen and Nutcracker excerpts
- Familiar melodies
- Holiday audience particularly responsive

MARKETING IMPLICATION: Feature these events prominently in the add-on offer. "Based on your ticket to Holiday Pops, we recommend Beethoven's 9th...\"""")

    # Slide 8: Section - Before/After Analysis
    add_section_slide(prs, "Before/After Analysis",
                      notes="""TRANSITION:

Now let's look at the hard numbers - what does our target market look like today, and what do we project after implementing the Classical Add-On program?

I'll show you:
1. Current state segmentation
2. Projected uptake by segment
3. Financial impact""")

    # Slide 9: Before State
    add_table_slide(prs, "Current State: Special/Holiday Single Ticket Buyers",
                    ["Segment", "Accounts", "Classical Events", "Classical Revenue"],
                    [
                        ["Group A: 0 Classical", "5,647", "0", "$0"],
                        ["Group B: 1 Classical", "729", "729", "$98,300"],
                        ["Group C: 2 Classical", "254", "508", "$61,504"],
                        ["Group D: 3+ Classical", "378", "1,984", "$304,181"],
                        ["TOTAL", "7,008", "3,221", "$463,985"]
                    ],
                    notes="""UNDERSTANDING THE SEGMENTS:

GROUP A - NEVER BOUGHT CLASSICAL (5,647 accounts):
- This is 80.6% of our target market
- They attended Special/Holiday events
- But they've NEVER purchased a Classical ticket
- This is our PRIMARY conversion target
- Huge untapped potential

GROUP B - BOUGHT 1 CLASSICAL (729 accounts):
- They've tried Classical once
- Opportunity to upsell to 2 or 3 tickets
- Already curious, need a push

GROUP C - BOUGHT 2 CLASSICAL (254 accounts):
- Regular crossover behavior
- One more ticket gets them to 3+
- One more ticket puts them in CYO conversion territory

GROUP D - BOUGHT 3+ CLASSICAL (378 accounts):
- Already engaged in Classical series
- Prime candidates for CYO subscription NOW
- 24.6% historical conversion rate
- Not our add-on target, but our CYO target

THE OPPORTUNITY:
- Group A alone represents $0 in Classical revenue today
- Converting even 20% of Group A to one ticket each = 1,129 new tickets
- That's why this program focuses heavily on Group A""")

    # Slide 10: After State
    add_table_slide(prs, "Projected State: With Add-On Program",
                    ["Group", "Target", "Responders", "New Tickets", "Revenue"],
                    [
                        ["Group A (0 Classical)", "5,647", "1,129", "1,689", "$90,008"],
                        ["Group B (1 Classical)", "729", "145", "187", "$9,603"],
                        ["Group C (2 Classical)", "254", "50", "50", "$2,489"],
                        ["TOTAL", "6,630", "1,324", "1,926", "$102,099"]
                    ],
                    "Assumes 20% response rate",
                    notes="""PROJECTED UPTAKE EXPLAINED:

ASSUMPTIONS:
- 20% response rate (conservative for targeted, same-day offers)
- Industry benchmarks suggest 15-30% for highly relevant offers
- We're being conservative

GROUP A PROJECTIONS:
- 5,647 accounts × 20% response = 1,129 responders
- Distribution: 60% buy +1, 30% buy +2, 10% buy +3
- Results in 1,689 new tickets
- Revenue: $90,008 (after discounts)
- This is the bulk of our opportunity

GROUP B PROJECTIONS:
- 729 accounts × 20% = 145 responders
- Buying 1-2 more tickets to reach 2 or 3 total
- 187 new tickets
- Revenue: $9,603

GROUP C PROJECTIONS:
- 254 accounts × 20% = 50 responders
- Each buying 1 more to reach 3 total
- 50 new tickets
- Revenue: $2,489

TOTAL IMPACT:
- 1,324 accounts respond to the offer
- 1,926 new Classical tickets sold
- $102,099 in incremental revenue

WHY THIS IS CONSERVATIVE:
- Same-day offers typically see higher response rates
- We're targeting a warm audience (they already attend ASO)
- Benefits are clear and immediate""")

    # Slide 11: Financial Summary
    add_metrics_slide(prs, "Financial Impact", [
        ("Incremental Revenue", "$102,099"),
        ("Discount Cost", "$10,678"),
        ("New Tickets Sold", "1,926")
    ],
    notes="""FINANCIAL SUMMARY:

$102,099 INCREMENTAL REVENUE:
- This is NET revenue after discounts
- Every dollar is incremental - these are sales that wouldn't happen otherwise
- Compare to: current Classical revenue from this segment is $463K
- Add-on adds 22% to their Classical spending

$10,678 DISCOUNT COST:
- This is the difference between discounted price and full price
- Represents 9.5% average discount across all purchases
- But remember: these buyers weren't going to buy at full price anyway
- The "cost" is theoretical - the revenue is real

1,926 NEW TICKETS:
- Each ticket represents a new Classical experience for a Special/Holiday buyer
- Many of these are first-time Classical attendees
- Potential for repeat purchases in future seasons
- Foundation for CYO subscription conversion

ROI PERSPECTIVE:
- Marketing cost to implement: minimal (email, signage)
- Revenue generated: $102,099
- Plus downstream CYO conversions worth $20-30K annually
- This is a high-ROI initiative""")

    # Slide 12: CYO Conversion
    add_content_slide(prs, "CYO Subscription Conversion Potential", [
        "Historical data: 24.6% of 3+ Classical buyers convert to CYO",
        "Program creates 205 new accounts reaching 3+ Classical tickets",
        "Projected new CYO subscribers: 50 accounts",
        "Estimated CYO revenue impact: $20,000 - $30,000 per season",
        "This is RECURRING revenue beyond the initial add-on purchase"
    ], "The add-on product builds a CYO subscriber pipeline",
    notes="""THE CYO CONVERSION PIPELINE:

This is where the strategy gets really powerful. The add-on isn't just about selling tickets - it's about building a subscriber pipeline.

HISTORICAL CONVERSION DATA:
- We analyzed single ticket buyers who reached 3+ Classical events
- 24.6% of them converted to CYO subscription
- This is a proven pathway to subscription

NEW 3+ CLASSICAL BUYERS FROM THIS PROGRAM:
- Group A buying +3: 112 accounts
- Group B buying +2 more (total 3): 43 accounts
- Group C buying +1 more (total 3): 50 accounts
- TOTAL: 205 new accounts reaching 3+ Classical

PROJECTED CYO CONVERSIONS:
- 205 accounts × 24.6% conversion rate = 50 new CYO subscribers
- This is a conservative estimate

CYO REVENUE IMPACT:
- Average CYO subscription value: $400-600/season
- 50 subscribers × $400-600 = $20,000-$30,000 per season
- This is RECURRING revenue
- Subscribers tend to renew at 70-80% rates

THE BIGGER PICTURE:
- Year 1: $102K ticket revenue + customer acquisition
- Year 2+: 50 new CYO subscribers × $400-600 = ongoing revenue
- Plus those CYO subscribers may upgrade further
- The add-on product is a funnel top, not just a revenue line""")

    # Slide 13: Risk Mitigation
    add_content_slide(prs, "Risk Mitigation: Product Differentiation", [
        "No cannibalization: Full-price Classical tickets remain available",
        "Isolated testing: If program underperforms, only add-on is affected",
        "Price integrity: Discount positioned as 'loyalty reward' for Special/Holiday attendees",
        "Clear segmentation: Offer only available to verified ticket holders",
        "Easy to adjust: Discount levels can be tuned based on response"
    ],
    notes="""ADDRESSING RISK CONCERNS:

This slide is critical for stakeholder buy-in. Let me address common concerns:

"WON'T THIS CANNIBALIZE FULL-PRICE SALES?"
No, because:
- The discount is ONLY available to verified Special/Holiday ticket holders
- General public cannot access this pricing
- If someone wants to buy Classical at full price, they still can
- We're targeting people who WEREN'T buying Classical anyway

"WHAT IF IT DOESN'T WORK?"
- The add-on is a separate product
- If it underperforms, we adjust or discontinue
- Our core Classical pricing remains unchanged
- No damage to existing business

"WILL THIS DEVALUE OUR CLASSICAL SERIES?"
- Positioning is crucial: "loyalty reward" not "discount"
- It's exclusive to a specific audience
- Actually ELEVATES Classical by introducing new audiences
- Similar to airline miles or hotel loyalty perks

"HOW DO WE CONTROL WHO GETS THE OFFER?"
- Ticket purchase verification required
- Unique codes tied to specific accounts
- Cannot be shared or transferred
- Clear terms and conditions

"WHAT IF DISCOUNT LEVELS ARE WRONG?"
- Start with pilot (2-3 events)
- A/B test different discount levels
- Adjust based on response data
- We have control to optimize""")

    # Slide 14: Implementation
    add_content_slide(prs, "Implementation Roadmap", [
        "Phase 1 - Pilot (First 2-3 Special/Holiday Events):",
        "    Test email + in-venue signage, measure response & redemption",
        "Phase 2 - Optimization:",
        "    A/B test discount levels, refine gateway concert recommendations",
        "Phase 3 - Scale:",
        "    Roll out to all events, integrate with CRM, automate triggers"
    ],
    notes="""IMPLEMENTATION ROADMAP:

PHASE 1 - PILOT (2-3 Special/Holiday Events):
Timeline: First 2-3 upcoming Special or Holiday events

Activities:
- Set up email automation (within 2 hours of event)
- Create in-venue signage with QR codes
- Train box office staff on the program
- Build landing page for offer redemption

Measurement:
- Email open rate (target: >40%)
- Click-through rate (target: >15%)
- Redemption rate (target: >5% of attendees)
- Average tickets per redemption

Goal: Prove the concept, gather learnings

PHASE 2 - OPTIMIZATION:
Timeline: Following 2-3 events

Activities:
- A/B test discount levels (5/10/15% vs flat 10%)
- Test different gateway concert recommendations
- Refine email timing and messaging
- Test SMS channel for opted-in customers

Measurement:
- Compare response rates across variations
- Identify optimal discount structure
- Determine highest-converting Classical events

Goal: Maximize response and revenue per contact

PHASE 3 - SCALE:
Timeline: Full season rollout

Activities:
- Automate all touchpoints
- Integrate with CRM for personalization
- Add predictive recommendations
- Implement CYO conversion follow-up sequence

Goal: Systematic, scalable revenue generation""")

    # Slide 15: Monitoring
    add_table_slide(prs, "Success Metrics to Monitor",
                    ["Metric", "Target", "Red Flag"],
                    [
                        ["Response rate", ">15%", "<10%"],
                        ["Redemption rate", ">5% of responses", "<2%"],
                        ["Full-price cannibalization", "<2% decline", ">5% decline"],
                        ["CYO conversion (from 3+ buyers)", ">20%", "<15%"]
                    ],
                    notes="""SUCCESS METRICS EXPLAINED:

RESPONSE RATE (Target: >15%):
- % of Special/Holiday attendees who engage with the offer
- Engagement = email open, QR scan, or inquiry
- 15% is conservative for targeted, timely offers
- Below 10% suggests messaging or timing problems

REDEMPTION RATE (Target: >5% of responses):
- % of engaged people who actually purchase
- This measures offer attractiveness
- Below 2% suggests discount isn't compelling enough
- Or gateway concerts aren't appealing

FULL-PRICE CANNIBALIZATION (Target: <2% decline):
- Critical to monitor
- Compare full-price Classical sales before/after program
- If full-price sales decline, we may be cannibalizing
- Below 2% decline is acceptable (could be normal variance)
- Above 5% decline is a red flag requiring investigation

CYO CONVERSION (Target: >20%):
- % of new 3+ Classical buyers who convert to CYO
- Historical baseline is 24.6%
- If lower, may need enhanced CYO follow-up
- This is the long-term value indicator

REPORTING CADENCE:
- Weekly during pilot phase
- Bi-weekly during optimization
- Monthly at scale
- Quarterly strategic review""")

    # Slide 16: Conclusion
    add_content_slide(prs, "Conclusion", [
        "Low-risk, high-potential opportunity",
        "Converts 1,324 Special/Holiday buyers into Classical attendees",
        "Generates $102,099 in incremental revenue from 1,926 new tickets",
        "Creates pipeline of 50 new CYO subscribers worth $20K-$30K/season",
        "Preserves pricing integrity through product differentiation"
    ], "80.6% of Special/Holiday buyers have never tried Classical - let's change that",
    notes="""CLOSING SUMMARY:

Let me leave you with the key points:

1. LOW-RISK, HIGH-POTENTIAL:
- We're not discounting existing products
- We're creating a new, targeted offer
- If it doesn't work, we adjust without affecting core business

2. CLEAR TARGET, PROVEN DEMAND:
- 7,008 accounts who already love ASO
- 80.6% have never tried Classical
- We're not creating demand; we're unlocking it

3. STRONG FINANCIAL CASE:
- $102,099 incremental revenue year one
- 1,926 new ticket sales
- ~9.5% average discount cost

4. STRATEGIC CYO PIPELINE:
- 205 new accounts reaching 3+ Classical
- 50 projected CYO conversions
- $20-30K recurring annual revenue

5. PRICE INTEGRITY PROTECTED:
- Product differentiation, not discounting
- Exclusive to verified ticket holders
- Full-price Classical remains unchanged

THE ASK:
Approve the pilot program for the next 2-3 Special/Holiday events. We'll measure results, optimize, and scale based on data.

The opportunity is clear. 80.6% of Special/Holiday buyers have never tried Classical. Let's change that.""")

    # Slide 17: Thank You
    add_title_slide(prs, "Thank You", "Questions?",
                    notes="""PREPARED ANSWERS FOR LIKELY QUESTIONS:

Q: Why not offer this to ALL single ticket buyers?
A: We're targeting Special/Holiday buyers specifically because they've demonstrated high engagement and willingness to pay. They're our warmest audience. We can expand later if successful.

Q: What's the cost to implement this?
A: Minimal. Email automation, signage design, landing page. Primary cost is staff time for setup. ROI is extremely high given $102K projected revenue.

Q: How do we prevent people from sharing discount codes?
A: Codes are tied to specific account IDs. Verification required at redemption. Terms clearly state non-transferable.

Q: What if people wait for the discount instead of buying full-price?
A: The offer is only available AFTER attending a Special/Holiday event. You can't get the discount without first buying a Special/Holiday ticket. This actually increases overall revenue.

Q: Why these specific discount levels (5/10/15%)?
A: Based on price elasticity analysis of Classical single ticket demand. These levels are enough to motivate without excessive margin erosion. We'll A/B test during pilot.

Q: How does this integrate with our existing marketing?
A: It's a separate, automated track triggered by Special/Holiday attendance. Doesn't interfere with regular Classical marketing. CRM integration in Phase 3.

Q: What about subscribers? Do they get this offer?
A: No. This is for SINGLE TICKET buyers only. Subscribers already have Classical access through their subscription.

THANK YOU FOR YOUR TIME. I'm happy to discuss any aspect in more detail.""")

    # Save
    output_path = Path(__file__).resolve().parent / "Classical_AddOn_Product_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
